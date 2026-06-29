#!/usr/bin/env python
"""Build the ~35-min program talk (mixed-researcher audience) in TWO formats:
  - presentation/haptic_brain_talk.pptx  (editable; PowerPoint/Keynote, or Google Slides)
  - presentation/haptic_brain_talk.pdf   (renders images everywhere)

Structure: Part 1 vision -> Part 2 literature & parameters -> Part 3 current study
-> Part 4 limits & next. Register: scientifically literate, not
neuroscience-specialist; lead with precise terms, analogies as intuition, and keep
the 'mechanism != proven' discipline the lit review insists on. Narration in .pptx notes.
Vision/lit content grounded in ~/Documents/LitretureReview_Github (papers by handle).
"""
from __future__ import annotations
from pathlib import Path
import textwrap

OUT = Path("presentation"); OUT.mkdir(exist_ok=True)
F = "results/dec4"
CF = "presentation/concept_figs"
TTL = "analysis/outputs/dec3/ttl_diagnostic/ttl_cannot_recover_tactor_phase.png"
TRIAL = "results/dec3/13_Teaching_and_Methods/trial_window_diagram.png"
SESSION_TL = "results/dec3/01_Session_Timeline/ttl_lfp_context_and_trials.png"
LOGO_W = "presentation/som_logo_white.png"       # Stanford School of Medicine, white
MOUSE = "presentation/concept_figs/mouse_white.png"
EVIDENCE = "presentation/concept_figs/evidence_ladder.png"
NAVY, TEAL, GREY, PANEL = "#2E2D29", "#8C1515", "#53565A", "#DAD7CB"   # Stanford: Black, Cardinal, Cool Grey, Fog
STONE, STONE_LIGHT, STONE_DARK, FOG_LIGHT, FOG_DARK = "#7F7776", "#D4D1D1", "#544948", "#F4F4F4", "#B6B1A9"

SLIDES = [
    dict(kind="title", title="Can vibrotactile stimulation drive brain clearance?",
         sub="Wearable vibrotactile stimulation, brain rhythms, and glymphatic clearance",
         sub2="vision and a pilot electrophysiology in a mouse",
         footer="a Snyder Lab (Stanford) × Buzsáki Lab (NYU) collaboration",
         footer2="Pardis Miri, PhD  ·  June 30, 2026",
         notes="The program question: can a wearable, vibrotactile stimulus engage brain rhythms in a way that helps the brain clear waste (the glymphatic system), with Alzheimer's relevance? I'll motivate it from the literature, lay out what to test, then show a first electrophysiology study testing the premise. I separate 'plausible mechanism' from 'proven' throughout. ~35 min — please interrupt."),

    # ===================== PART 1 — VISION =====================
    dict(kind="section", title="Part 1 · The vision — clearing the brain with vibrotactile stimulation",
         sub="Where clearance biology could matter: Alzheimer's amyloid/tau · stroke recovery/inflammation",
         refs="Refs: Iliff 2012; Xie 2013; Murdock 2024 · Gaberel 2014; Mestre 2020; JCI 2026",
         sub_y=5.35,
         notes="Disease examples for the vision slide. Keep this as mechanism motivation, not a treatment claim. Alzheimer's: the glymphatic idea starts with clearance of interstitial solutes including amyloid-beta (Iliff 2012), sleep-dependent metabolite clearance (Xie 2013), and 40 Hz multisensory stimulation increasing glymphatic Aβ clearance in mouse models (Murdock 2024). Stroke: glymphatic perfusion is impaired after several stroke models (Gaberel et al., Stroke 2014); acutely, glymphatic CSF influx can also worsen edema (Mestre et al., Science 2020), so timing matters; a new mouse chronotherapy study reports increased glymphatic function, reduced cytokine burden, and improved post-stroke recovery (JCI 2026). Bottom line to say out loud: if a wearable can shape brain state safely, the broader long-term application space may include diseases where clearance/inflammation/protein propagation matter, but our present data only test neural target engagement."),
    dict(kind="content", kicker="the problem", title="Alzheimer's resists drugs — non-invasive neuromodulation is an open frontier", fig=f"{CF}/landscape.png",
         take="Pharmacology is real but incomplete: modest slowing with BBB/safety burden. Multidomain lifestyle has the largest human-scale evidence; Hosseini's Stanford work adds a relevant cognitive-training branch. Neuromodulation is promising but still early.",
         notes=("Framing the gap. The three boxes are the landscape in one glance. Pharmacology: do not say 'dead end' - say 'incomplete.' Anti-amyloid antibodies now show real but modest slowing, and the blood-brain barrier plus safety monitoring still make this a hard drug-delivery problem. Lecanemab evidence: van Dyck et al., 'Lecanemab in Early Alzheimer's Disease,' New England Journal of Medicine, 2023; about 27% slower clinical decline over 18 months, with infusion, amyloid-confirmation, MRI, and ARIA monitoring burden. Donanemab evidence: Sims et al., 'Donanemab in Early Symptomatic Alzheimer Disease: The TRAILBLAZER-ALZ 2 Randomized Clinical Trial,' JAMA, 2023; about 29-36% slowing depending on outcome/population, also with ARIA and monitoring burden. "
                "Lifestyle/multidomain: 'multidomain' means several risk pathways are targeted together, not one habit alone. The domains include exercise, diet, cognitive challenge/training, social engagement, and cardiovascular/metabolic monitoring. FINGER evidence: Ngandu et al., 'A 2 year multidomain intervention of diet, exercise, cognitive training, and vascular risk monitoring versus control to prevent cognitive decline in at-risk elderly people (FINGER): a randomised controlled trial,' The Lancet, 2015; NTB total score improved 0.20 vs 0.16 over 2 years, a modest but positive signal. U.S. POINTER evidence: Baker et al., 'Structured vs Self-Guided Multidomain Lifestyle Interventions for Global Cognitive Function: The US POINTER Randomized Clinical Trial,' JAMA, 2025; structured improved 0.243 SD/year vs 0.213 SD/year for self-guided, difference +0.029 SD/year, p=.008. "
                "Hadi Hosseini belongs here, but carefully: his work is not the same broad lifestyle package as FINGER/U.S. POINTER. It is a Stanford cognitive-training / neurocognitive-rehabilitation branch of the multidomain prevention space. Gozdas et al., including S. M. Hadi Hosseini, 'Long-term cognitive training enhances fluid cognition and brain connectivity in individuals with MCI,' Translational Psychiatry, 2024; 6-month multi-domain computerized cognitive training in amnestic MCI, 34 completers; Fluid Cognition improved 10.9% in the cognitive-training group, Cohen's d=0.70, but the sample is small, so call it promising rather than definitive. Related Hosseini background I checked but did not make the main impact anchor: Hosseini, Kramer, and Kesler, 'Neural correlates of cognitive intervention in persons at risk of developing Alzheimer's disease,' Frontiers in Aging Neuroscience, 2014, a review/methods framing paper; and Bruno, Shaw, and Hosseini, 'Toward Personalized Cognitive Training in Older Adults: A Pilot Investigation of the Effects of Baseline Performance and Age on Cognitive Training Outcomes,' Journal of Alzheimer's Disease, 2024 volume / online first 2023, a small healthy-aging pilot about personalization. Single-domain caution: Barnes et al., 'Trial of the MIND Diet for Prevention of Cognitive Decline in Older Persons,' New England Journal of Medicine, 2023; MIND diet plus caloric restriction did not significantly beat control diet with caloric restriction, difference 0.035 standardized units, p=.23. "
                "Non-invasive neuromodulation: latest human methods include personalized precuneus rTMS, home gamma tACS, and 40 Hz audiovisual sensory stimulation. rTMS evidence: Koch et al., 'Effects of 52 weeks of precuneus rTMS in Alzheimer's disease patients: a randomized trial,' Alzheimer's Research & Therapy, 2025; CDR-SB worsened +1.36 rTMS vs +2.45 sham at 52 weeks, but only 48 enrolled and 31-32 reached the long follow-up. Gamma tACS evidence: Cantoni et al., 'Home-Based Gamma Transcranial Alternating Current Stimulation in Patients With Alzheimer Disease: A Randomized Clinical Trial,' JAMA Network Open, 2025; n=50, home gamma tACS showed small clinical differences and EEG engagement, but this is still proof-of-concept. GENUS human sensory evidence: Chan et al., 'Gamma frequency sensory stimulation in mild probable Alzheimer's dementia patients: Results of feasibility and pilot studies,' PLOS ONE, 2022; feasibility/pilot studies suggest safety and target engagement, not definitive efficacy. Bottom line to say out loud: lifestyle has the strongest human-scale evidence but modest effects; Stanford/Hosseini supports the cognitive-training branch; neuromodulation is exciting because it may be targetable/wearable, but it still has to prove durable clinical and biomarker impact.")),
    dict(kind="content", kicker="why clearance", title="A leading hypothesis: failing glymphatic clearance lets amyloid build up", fig="presentation/cells-13-00286-g001.png",
         fig_maxw=9.6,
         fig_box=(0.14, 0.24, 0.72, 0.48),
         take="Simple idea: sleep boosts CSF–brain-fluid exchange, helping waste clearance. If clearance fails, amyloid-β / tau may build up — an active but still-unsettled AD hypothesis.   Figure: Szlufik et al., Cells 2024;13:286 (CC BY).",
         notes="The mechanistic bridge — why stimulating the brain might help at all. One leading hypothesis is that AD is partly a clearance failure. The glymphatic system (Iliff 2012) moves CSF along perivascular spaces to wash amyloid-β and tau out toward the cervical lymph nodes; it is strongly boosted during sleep (Xie 2013) and by vascular pulsation, and depends on the AQP4 water channel. In AD, amyloid is cleared more slowly (Mawuenyega 2010), and impaired glymphatic flow on DTI-ALPS imaging appears to precede amyloid deposition (ADNI 2024). Important caveat for this audience: it remains contested — DTI-ALPS may index vascular rather than glymphatic function (Mayo 2025), and some studies find clearance reduced rather than increased in sleep. So present it as an active, plausible contributor, not the proven sole cause."),
    dict(kind="content", kicker="why rhythms", title="Neural waves can organize CSF perfusion and clearance", fig=f"{CF}/csf_wave_evidence.png",
         take="Jiang-Xie et al. 2024 (Nature) is the cleanest causal anchor: flatten neural/ionic waves and CSF infiltration/clearance drops; generate synthetic waves with transcranial optogenetics and CSF-to-ISF perfusion increases. Sleep, neuromodulator, vascular, and 40 Hz sensory papers support the broader state-control story.",
         notes=("This is the dedicated mechanism slide. Say it simply: the glymphatic pathway is not just plumbing; brain state and brain waves can organize the flow. The strongest causal paper is Jiang-Xie et al., 'Neuronal dynamics direct cerebrospinal fluid perfusion and brain clearance,' Nature, 2024. Their method was not vibrotactile stimulation. They measured large rhythmic ionic waves in brain tissue, then did two causal manipulations: chemogenetically flattening those waves impaired CSF infiltration and molecular clearance, while synthetic waves generated by transcranial optogenetic stimulation increased CSF-to-interstitial-fluid perfusion. That is why I call it the causal anchor. "
                "Then make the evidence ladder clear. Fultz et al., 'Coupled electrophysiological, hemodynamic, and cerebrospinal fluid oscillations in human sleep,' Science, 2019, showed in humans that NREM slow neural waves, blood oxygenation waves, and CSF flow oscillations are coupled. Hablitz et al., 'Increased glymphatic influx is correlated with high EEG delta power and low heart rate in mice under anesthesia,' Science Advances, 2019, showed glymphatic influx tracks slow cortical activity/state in mice. Hauglund et al., 'Norepinephrine-mediated slow vasomotion drives glymphatic clearance during sleep,' Cell, 2025, adds a sleep-pump mechanism: norepinephrine-linked slow vasomotion during NREM sleep drives clearance. Chuang et al., 'Cholinergic basal forebrain neurons regulate vascular dynamics and cerebrospinal fluid flux,' Nature Communications, 2025, shows neuromodulatory control of BOLD-CSF coupling and glymphatic flux, with human imaging plus mouse lesion/pharmacology. Broggini et al., 'Long-wavelength traveling waves of vasomotion modulate the perfusion of cortex,' Neuron, 2024, supports the vascular-wave side: vasomotion waves modulate cortical perfusion, though it is more blood/perfusion than direct glymphatic clearance. Murdock et al., 'Multisensory gamma stimulation promotes glymphatic clearance of amyloid,' Nature, 2024, connects external 40 Hz light+sound to increased CSF influx/ISF efflux and amyloid clearance in AD-model mice. Bottom line: Jiang-Xie is the strongest direct wave-to-clearance causal paper; the others support the broader idea that neural state, neuromodulators, and vascular waves regulate CSF/glymphatic flow. This motivates our open question, but it does not prove vibrotactile stimulation will do the same.")),
    dict(kind="content", kicker="and it can be driven", title="And clearance can be driven non-invasively — preclinical evidence", fig=f"{CF}/driving_clearance.png",
         take="This slide is animal proof-of-principle, not human clinical proof. 40 Hz sensory stimulation increased glymphatic flow and reduced amyloid in AD-model mice (Murdock 2024). Focused ultrasound increased CSF influx/clearance in rodents (Aryal 2022; Azadian 2025). Open question: can vibrotactile stimulation enhance glymphatic clearance?",
         notes="The proof-of-principle that motivates the program. Be very clear: this slide is mostly preclinical animal evidence, not human efficacy. The previous slide made the mechanism point that neural and vascular dynamics can organize CSF perfusion. This slide asks: can an outside intervention move that biology? Murdock et al. 2024, Nature, gives the directly relevant sensory-stimulation result: 40 Hz multisensory light + sound in 5XFAD AD-model mice increased CSF tracer influx, increased ISF efflux, lowered frontal-cortex amyloid by about 30%, and blocking the AQP4 water channel abolished the amyloid-clearance effect, tying it to the glymphatic route. Separately, Stanford's focused-ultrasound work is also preclinical: Aryal 2022 used transcranial ultrasound to increase CSF influx and intrathecal drug delivery in rats; Azadian 2025 used low-intensity focused ultrasound to clear debris and improve outcomes in mouse hemorrhagic-stroke models. Together they do not prove a wearable haptic treatment works in humans. They support the narrower claim: clearance biology can be modulated from outside the skull in animal models, so it is reasonable to ask whether vibrotactile stimulation can enhance glymphatic clearance."),
    dict(kind="content", kicker="the thesis", title="The proposal: use vibrotactile stimulation to engage brain rhythms and drive clearance", fig=f"{CF}/vision_chain.png",
         take="Wearable vibrotactile stimulation → entrain brain rhythms → engage glymphatic clearance of amyloid-β / tau. A mechanism hypothesis — explicitly not yet a treatment claim.",
         notes="The core chain. The glymphatic system is the brain's waste-clearance pathway — CSF/interstitial-fluid exchange that removes amyloid-β and tau. The hypothesis: if we can drive the right rhythms or brain-states from the skin, we may enhance that clearance. This is a mechanism hypothesis to test, not a clinical claim — I'll be explicit about that distinction the whole way through."),
    dict(kind="content", kicker="why haptics", title="Why haptics, and why wearable", fig=f"{CF}/whitespace.png",
         take="Cognito Spectris/GammaSense: Alzheimer's trials using audiovisual 40 Hz. Us: a 2-finger fingertip vibrotactile glove designed for long sessions, sleep-compatible use, and pairing with wearable sensing.",
         notes="The niche. The established sensory-stimulation approach is audiovisual 40 Hz (GENUS), now represented clinically by Cognito's Spectris / GammaSense device program for Alzheimer's disease, mostly mild-to-moderate AD populations. I also include Mlinaric et al. 2025 as an independent human iEEG target-engagement result, not a Tsai/Cognito clinical study and not an AD efficacy trial: 40 Hz visual stimulation produced 40 Hz responses beyond visual cortex, including hippocampal contacts, in 11 epilepsy patients with 490 contacts. These are device studies, so ClinicalTrials.gov lists the phase as N/A; the more useful program stages are feasibility, pivotal, open-label extension, and dose-ranging/safety. OVERTURE (NCT03556280, n=76 randomized in the peer-reviewed paper; mild-to-moderate AD) was the feasibility study: safe and tolerable, but the primary efficacy endpoint did not separate active from sham; exploratory secondary outcomes favored active treatment on ADCS-ADL, MMSE, and MRI whole-brain-volume loss. HOPE (NCT05637801, registry enrollment n=670; mild-to-moderate AD) is the pivotal trial. The OLE HOPE study (NCT06245031, registry enrollment n=402) is the open-label active extension: people who complete the blinded HOPE trial can enter a follow-up where everyone receives the active device, with no sham/control group and no blinding, mainly to collect longer-term safety, adherence, and durability data. They are still ongoing/no posted results and are testing whether daily 40 Hz sensory stimulation slows 12-month clinical progression using ADCS-ADL/MMSE-style function-cognition endpoints. ETUDE (NCT03661034, registry enrollment n=20) was a dose-ranging/safety/amyloid PET study in cognitive impairment / AD-relevant participants, but I found no posted results. In parallel to audiovisual stimulation, our route is haptics: a vibrotactile device applied to the fingertips, formed as a two-finger glove for the index and middle fingers, designed for long sessions, sleep-compatible use, and pairing with wearable sensing. The human testing path is intentionally resource- and funding-driven: possible routes include Stanford Hospital iEEG patients for direct target-engagement physiology, stroke patients, or Alzheimer's disease patients, depending on the clinical collaboration and funding path that opens first."),

    # ===================== PART 2 — LITERATURE & PARAMETERS =====================
    dict(kind="section", title="Part 2 · What the literature says — and what to test"),
    dict(kind="content", kicker="evidence · 40 Hz gamma",
         title="The 40 Hz evidence chain: from mouse GENUS to early human extension",
         fig=f"{CF}/genus_evidence_chain.png",
         take="GENUS moved from mouse pathology reduction to an AQP4/glymphatic mechanism and early human feasibility signals; controlled human efficacy remains the open question.",
         notes="The anchor evidence chain. The GENUS/Tsai series began by showing that 40 Hz visual/auditory stimulation can entrain gamma activity and reduce amyloid/tau pathology in AD-model mice; Murdock 2024 then linked 40 Hz stimulation to AQP4-dependent glymphatic CSF influx and amyloid clearance; Park & Tsai 2025 summarized the expanding preclinical and early clinical evidence base; and Chan et al. 2025 reported that a small open-label human extension of daily 40 Hz GENUS was feasible/safe and suggested possible longer-term cognitive/biomarker benefit, but without a blinded control group. This is why 40 Hz is the canonical frequency to start from, while still separating target engagement and mechanism from definitive human efficacy."),
    dict(kind="content", kicker="vocabulary before caveat",
         title="Vocabulary first: from 'the brain noticed' to 'the rhythm captured it'",
         fig=EVIDENCE, full_fig=True,
         notes="Transition slide before the Buzsáki/Soula caveat. Say: before I explain the caveat, we need a shared vocabulary for what counts as a weak versus strong stimulation claim. One clap makes the body react = evoked response. Different clap speeds or strengths producing different neural changes = stimulus-specific or dose-response. A steady sequence appearing as a spectral peak = steady-state or frequency-following, but that alone can be LFP pickup or a sensory-following response. Phase locking means activity lands at the same point in each stimulus cycle, which requires knowing the actual delivered vibration phase. Entrainment is the strongest claim: the outside rhythm captures or pulls an endogenous rhythm. This ladder is why the next slide is not semantic nitpicking — it defines what evidence is required before using the word entrainment."),
    dict(kind="content", kicker="the key caveat",
         title="Steady-state response is not entrainment: the Buzsáki/Soula guardrail",
         fig=f"{CF}/response_entrainment_guardrail.png",
         take="A driven 40 Hz line can be a real sensory response without proving that native gamma was captured. We therefore separate response, frequency-following, phase-locking, and true entrainment.",
         notes=("This is the credibility slide. The local Buzsáki email inventory is marked confidential, so do not quote the email or the attached critique; paraphrase only. Core paraphrase: sensory steady-state responses are easy to evoke across modalities and frequencies; resonance and attention can amplify a 40 Hz line; that still does not prove the brain's own gamma generator has been entrained. The local inventory lists Herrmann 2001 as the steady-state/flicker-resonance citation: human EEG responses to 1-100 Hz flicker show resonance phenomena in visual cortex. It lists Tiitinen et al. 1993, Nature, as the attention citation: selective attention enhances the human auditory 40 Hz transient response. "
                "Public anchor: Soula, Martín-Ávila, Zhang, Dhingra, Nitzan, Sadowski, Gan, and Buzsáki, Nature Neuroscience 2023, tested 40 Hz visual flicker in APP/PS1 and 5xFAD Alzheimer's mouse models using silicon probes in visual cortex, entorhinal cortex, and hippocampus. Their headline result was not 'nothing happened anywhere': visual cortex responded. The important gradient was visual cortex >> hippocampus / entorhinal cortex. Train-onset spike responses were V1 92/458 neurons (20.1%), EC 1.9% of 211 units, and CA1 hippocampus 0.2% of 463 units. Phase-locking to individual 40 Hz cycles was V1 28.8%, CA1 6.7%, and EC 3.3%. Population firing-rate changes during 40 Hz on/off were significant in V1 but not in CA1 (Extended Data Fig. 5, hippocampus p=0.2581). They also reported no reliable plaque, microglia, or Aβ40/42 reductions. "
                "Independent caution: Duecker, Gutteling, Herrmann, and Jensen 2021 in Journal of Neuroscience showed in human MEG that rapid visual flicker can produce a strong flicker response, but endogenous gamma and the flicker response can coexist in visual cortex without phase/frequency entrainment. Yang and Lai 2023 is another negative pathology replication in 5xFAD mice after chronic 24/40/80 Hz visual stimulation, while newer positive papers suggest protocol intensity and implementation may matter. "
                "So for our haptic program, the language discipline is: say response, target engagement, narrow-band response, or frequency-following unless we have delivered-stimulus phase, LFP/spike timing, phase/frequency capture of a native rhythm, specificity, and controls. This tees up the response-to-entrainment ladder and protects the project from overclaiming.")),
    dict(kind="content", kicker="evidence · clearance",
         title="Entrainment or not — did glymphatic flow change?",
         fig=f"{CF}/glymphatic_measurement_methods.png",
         take="Take-home: human CSF motion → Kipnis causal neural-wave perfusion → Murdock 40 Hz amyloid clearance; haptics needs flow/clearance biomarkers.",
         notes=("This slide is about measurement confidence, not just biological plausibility. First: human CSF-motion imaging. Fultz et al., Science 2019, used simultaneous EEG/fMRI during human NREM sleep and showed slow neural activity, hemodynamic/BOLD waves, and ventricular CSF inflow oscillations are coupled. Effect-size anchor: Fultz reported about +5.5 dB CSF power at ~0.05 Hz during NREM, and Williams et al., PLOS Biology 2023, reported +6-10% evoked CSF signal across human visual-stimulation datasets. Williams extended the idea to awake sensory stimulation: visual stimulation can drive large-scale CSF flow signals. Takeaway: human brain state and sensory-evoked hemodynamics can move macroscopic CSF; caveat: this is noninvasive and human-relevant, but it is still an indirect CSF-motion readout, not direct molecular clearance. "
                "Second: the Kipnis causal bridge. Jiang-Xie et al., Nature 2024, 'Neuronal dynamics direct cerebrospinal fluid perfusion and brain clearance,' measured neuronal/ionic waves and CSF tracer perfusion in mice. Effect-size anchor: sleep/wake tracer-flow readouts are roughly ~1.5x whole-slice and ~2x in unsilenced hippocampus; chemogenetic silencing impairs tracer entry and clearance. Flattening the waves reduced CSF infiltration and molecular clearance, while artificially generated transcranial optogenetic waves increased CSF-to-ISF perfusion. Takeaway: neural dynamics can direct fluid perfusion; caveat: this is mouse/preclinical and not AD- or 40 Hz-specific. "
                "Third: the 40 Hz clearance bridge. Murdock et al., Nature 2024, 'Multisensory gamma stimulation promotes glymphatic clearance of amyloid,' used cisterna magna tracer influx, cortical tracer imaging/two-photon style flow assays, ISF efflux measures, amyloid movement toward deep cervical lymph nodes, and AQP4 disruption. Effect-size anchor: 40 Hz audiovisual stimulation produced about ~4x CSF influx, ~1.5x ISF efflux, and about ~30% lower cortical amyloid signal. Takeaway: 40 Hz audiovisual GENUS increased glymphatic flow measures and AQP4-dependent amyloid clearance in AD-model mice; caveat: closest to our story, but still mouse evidence and audiovisual rather than haptic. "
                "Say the bottom line clearly: for the haptic program, entrainment is only a target-engagement question. The stronger program endpoint is whether vibrotactile brain-state modulation changes measurable CSF/glymphatic flow, drainage, or clearance biomarkers.")),
    dict(kind="content", kicker="evidence · vibrotactile",
         title="The vibrotactile anchor: promising AD-model evidence, but no direct entrainment test",
         fig=f"{CF}/vibrotactile_anchor_suk.png",
         take="Suk 2023 used whole-body 40 Hz vibration from a large speaker/platform and saw c-Fos/pathology signals in superficial SSp/MOp cortex. It did not directly record hippocampus, entorhinal cortex, LFP phase-locking, single units, or entrainment — the gap our study targets.",
         notes=("This is the vibrotactile-specific anchor, but state exactly what kind of anchor it is. Suk et al. 2023 used whole-body 40 Hz vibration in AD-model mice, delivered with a speaker/platform style setup, not a wearable fingertip device. Their brain readout was not implanted electrophysiology. They used immunostaining and related tissue/pathology readouts. The immediate activity evidence was increased c-Fos activity marker signal in superficial cortical regions, especially primary somatosensory cortex / SSp and primary motor cortex / MOp. "
                "The important limitation for our story: they did not directly record hippocampus or entorhinal cortex, and they did not test LFP phase-locking, single-unit responses, or entrainment to the vibration cycle. That does not invalidate the paper; it defines the next measurement. The slide should make the transition cleanly: Suk shows haptics can matter in an AD-model context, but our program has to prove direct target engagement in relevant circuits with electrophysiology, then eventually flow/clearance biomarkers.")),
    dict(kind="content", kicker="→ what to test", title="Which frequencies should we test, and why?",
         fig=f"{CF}/frequency_menu.png",
         take="Rationale, not one magic frequency: 40 Hz is the AD/GENUS/Suk comparator; 50 Hz is our data-driven candidate; low/mid bands test specificity. Slow/patterned rows are optional mechanistic extensions.",
         notes=("The payoff of Part 2. This slide turns the literature into an experimental menu. The minimum defensible haptic panel is 5, 10, 20, 26, 40, and 50 Hz. 5 and 10 Hz are low-flutter controls from Mountcastle/Romo and Hayashi-style tactile coding work; they ask whether any periodic touch is enough to drive the brain. 20 and 26 Hz are mid-frequency controls: 20 Hz is a standard comparator in the 40 Hz sensory-stimulation literature, and 26 Hz keeps continuity with our Dec 3/Dec 4 recordings, where it produced a large gross LFP transient but not the cleanest single-unit peak. "
                "40 Hz is required because it is the AD/GENUS anchor from Tsai/Singer/Martorell, the Murdock 2024 glymphatic bridge, and the Suk 2023 vibrotactile AD-model precedent. But because of the Buzsáki/Soula guardrail, phrase it as testing target engagement and frequency-specific response, not automatically entrainment. 50 Hz is required because our Dec 4 data show the strongest clean single-unit firing-rate modulation in dorsal hippocampus and lateral entorhinal cortex, and it also sits at the upper end of classic tactile flutter. "
                "Slow ~0.05-1 Hz is not a core haptic frequency here; it is a clearance-state hypothesis motivated by Fultz human sleep CSF waves and Jiang-Xie/Kipnis 2024 synthetic 1 Hz neural-wave perfusion. The important message is not that one number is magic; it is that each condition has a reason and a failure mode.")),

    # ===================== PART 3 — CURRENT STUDY =====================
    dict(kind="section", title="Part 3 · Does vibrotactile stimulation reach the brain? A first study",
         sub="Foundational electrophysiology in the mouse — before any clearance claim"),
    dict(kind="content", kicker="orientation · anatomy",
         title="Where our probes sit — and how a body signal could reach them",
         fig=f"{CF}/lec_dhpc_pathway.png",
         take="Our probes are in lateral entorhinal cortex (LEC) and dorsal hippocampus (dHPC) — deep medial-temporal targets, not first-touch cortex. A back-vibration signal reaches them only through a long, multisynaptic route (see diagram), many synapses past where body touch first arrives in S1.",
         notes="Orientation slide before the study design: it says where the electrodes are and why that shapes what we can expect. Our probes are NOT in first-touch cortex — they are in lateral entorhinal cortex (LEC) and dorsal hippocampus (dHPC), deep medial-temporal structures many synapses downstream of where body touch first arrives. The canonical route: the mouse back is hairy skin, so 5-50 Hz vibration is the tactile 'flutter' band carried by rapidly-adapting hair-follicle/field afferents (not Pacinian corpuscles — those are tuned to ~300-500 Hz and in rodents sit in distal-limb periosteum, not back skin). Those afferents enter via thoracic dorsal root ganglia, ascend the dorsal columns to the gracile/cuneate nuclei, cross, and relay through thalamic VPL to S1 — the trunk area, which is the FIRST cortical stage. "
                "From S1 the signal climbs the cortico-cortical hierarchy: S1 -> S2 / posterior parietal -> perirhinal/postrhinal cortex (the multimodal gateway that preferentially feeds LEC; Burwell & Amaral 1998) -> LEC (the parahippocampal multimodal 'what'/content integrator; Doan 2019, Deshmukh & Knierim 2011) -> via the perforant path into the dorsal hippocampus. The crucial physiology: primate S1 can phase-lock cycle-by-cycle to flutter, but that temporal code collapses about ten-fold at the very next synapse (S2), where frequency becomes a firing-RATE code (Salinas/Romo). LEC and dHPC are several synapses further still, so the prior is clear — expect stimulus-driven single-unit RATE modulation, but NOT clean 40-50 Hz phase-locked LFP entrainment. That is exactly the pattern in our data. "
                "Keep one caveat visible (dashed route): a whole-body buzz is also an arousing stimulus, so it can modulate LEC/dHPC through ascending arousal/neuromodulatory systems (locus coeruleus, cholinergic basal forebrain, reticular activation) as a STATE change, not a frequency-specific sensory drive — which at the single-unit-rate level is hard to distinguish from faithful sensory routing. Figure is a teaching schematic, not an atlas plate (Allen CCF; Kandel PNS; Burwell & Amaral 1998; Doan 2019; Witter)."),
    dict(kind="content", kicker="study design", title="The conditions we tested — and a real session for perspective", fig=SESSION_TL,
         video="presentation/experiment_video.MOV", video_poster="presentation/experiment_video_poster.png", video_size=(1080, 1920),
         take="We varied carrier frequency (5/10/26/50 Hz) × 3 amplitudes, in 3 s ON / 3 s OFF blocks, ~200 randomized repeats per condition — recording hippocampus and entorhinal cortex together. Left: the setup in action; right: a real Dec 3 session (15-min baseline → 1200 trials over 120 min → post).",
         notes="These are the conditions we tested: carrier frequency (5, 10, 26, 50 Hz) crossed with three amplitudes, delivered as 3 s ON / 3 s OFF blocks, ~200 randomized repeats per condition, while recording hippocampus and entorhinal cortex simultaneously with linear silicon probes. Rather than a schematic, this figure shows a real session for perspective. Top row: the accelerometer/vibration sensor across the whole ~3-hour recording — the 15-minute quiet baseline, then ~120 minutes of dense stimulation (1200 trials in randomized order), then the post-experiment period. Middle row: the simultaneously recorded hippocampal LFP over the same timeline. Bottom row: a zoom into a few consecutive trials — each a different randomly chosen condition (gold = 3 s ON/vibrating, grey = 3 s OFF/rest), with vibration start/end marks and the motion sensor. For Dec 3 the conditions were the 5 and 26 Hz carriers at three amplitudes (200 repeats each = 1200 trials); Dec 4 used the 10 and 50 Hz carriers. Each analysis window is referenced to the pre-stimulus baseline, with 100 ms margins to isolate onset/offset transients. Randomizing order guards against slow drift / arousal confounds. TTL alignment landed 1185/1200 stimulus pulses inside the ON window (median onset ~175 ms)."),
    dict(kind="content", kicker="key methodological distinction", title="One electrode yields two readouts: LFP and single-unit activity", fig=f"{CF}/stadium.png",
         take="LFP = low-frequency aggregate field (many cells). A single unit = one individual neuron's firing (spikes), isolated by spike sorting. They differ in what can fake them — spikes are the conservative readout.",
         notes="The crux of the whole study. An extracellular electrode captures two regimes. The local field potential (LFP) is the low-frequency aggregate of synaptic/transmembrane currents from many neurons — large, but non-specific and susceptible to volume-conducted or instrumental noise. A 'single unit' is one individual neuron's spiking, isolated by spike sorting (clustering spikes by waveform). Stadium intuition: the crowd's roar (LFP) versus one identified shout (a single unit). The asymmetry that matters: an electrical artifact can add power to the LFP, but it cannot make a sorted neuron fire — so single-unit rate changes are the conservative, artifact-resistant evidence."),
    dict(kind="content", title="Dec 4 spectrogram: the narrowband LFP effect appears at 50 Hz, mainly in LEC",
         fig=f"{F}/05_Frequency_Spectral/trial_avg_spectrogram_dec4.png",
         take="At 50 Hz, LEC shows an amplitude-graded narrowband line; dHPC does not. But an electromechanical actuator can inject a 50 Hz artifact into the LFP, so this line is artifact-suspect until proven otherwise.",
         notes="This is the matched Dec 4 spectrogram. Top row is LEC, bottom row is dHPC; columns are amp100, amp180, amp250 at 50 Hz. The LEC 50 Hz band grows with amplitude, while dHPC lacks a comparable 50 Hz band. Be disciplined in wording: this is a real measured LFP spectral line, but an electromechanical actuator and its drive electronics can couple a 50 Hz signal into the recording via wiring or volume conduction. In the LFP alone, that is indistinguishable from neural 50 Hz. So this slide should tee up the actual contamination evidence immediately, rather than making a separate cartoon/confound slide. The key phrase: measured LFP effect, artifact-suspect."),
    dict(kind="content", kicker="finding 3, controlled", title="Disconnected channels carry the 50 Hz — so it is largely pickup",
         fig=f"{F}/12_ChannelQC_Traces/50hz_pickup_gradient_dhpc_vs_lec.png",
         take="Disconnected electrodes (which cannot record neural signal) show ~6× more 50 Hz than in-tissue channels. The entorhinal 50 Hz LFP is substantially artifact.",
         notes="The discriminator: broken/disconnected channels cannot record neurons but still act as antennas for pickup. If they carry the 50 Hz, it's instrumental. They do — ~6× the in-tissue channels — so the entorhinal 50 Hz LFP is largely stimulator artifact. The headline LFP result is mostly the machine. This is the central twist, and exactly the 'steady-state response ≠ real engagement' caution from Part 2, caught in our own data."),
    dict(kind="content", kicker="the artifact, in full",
         title="…and the 50 Hz reaches a third of the good LEC channels too — not only the dead ones",
         fig=f"{F}/12_ChannelQC_Traces/fiftyhz_tissue_contamination_dec4.png",
         take="It is not just the disconnected electrodes: about a third of the tissue-good LEC channels also carry the 50 Hz pickup (hippocampal tissue: zero). So no single entorhinal LFP contact is trustworthy at 50 Hz — which is exactly why the single-unit readout, detected above 300 Hz, is the one to believe.",
         notes="Sharpening the artifact conclusion before pivoting to the clean result. The previous slide showed disconnected channels carry the 50 Hz. But the pickup is not confined to dead electrodes — through the shared reference and volume conduction it bleeds into the tissue-good channels too. Per channel: hippocampal tissue is clean (0% contaminated), about a third of the good entorhinal tissue channels are clearly contaminated, and the dead ones are far worse. The honest implication is that you cannot trust any single entorhinal LFP contact at 50 Hz. This is NOT a reason to distrust the spikes: spike detection is high-pass, above 300 Hz, so the 50 Hz pickup is removed before a spike is ever detected. It is the reason we lean on the single-unit readout, not the LFP — and note that all entorhinal units happen to sit on these contaminated channels, which is why the spike-level controls (coming up) matter."),
    dict(kind="content", kicker="finding 4 · unit by unit", title="The clean readout pickup can't fake: single units change firing at 50 Hz",
         fig=f"{F}/11_Spikes/raster_psth_examples_dec4.png",
         take="Four example single units (hippocampus & entorhinal, up- and down-modulated): each tick is one spike, each row one of 200 trials; the PSTH below averages them. Firing shifts during the 3 s ON window (red) and returns during OFF (blue), measured against a never-stimulated pre-study baseline (orange). By cell type, the hippocampal 50 Hz up-drive is interneuron-led, while entorhinal principal cells are suppressed.",
         notes="The per-unit picture behind the population result. Top of each column is a spike raster — one row per trial (200 trials), each vertical tick a single spike from that sorted neuron; below it the trial-averaged firing rate (PSTH). Red shade is the 3 s ON window, blue the 3 s OFF window. The orange dashed line is each unit's own firing rate during the 30 minutes of quiet recording before the study ever started — the never-stimulated reference — so you can see not just the ON-vs-OFF change but where the firing sits relative to baseline. Two units drive up (↑) during 50 Hz and two are suppressed (↓); all return toward baseline in OFF. This is the readout pickup cannot fake: an electrical artifact can add power to the field, but it cannot place or remove a sorted neuron's spikes."),
    dict(kind="content", kicker="finding 4 · all Dec 4 units",
         title="It holds across the whole population — every curated dHPC + LEC unit",
         fig=f"{F}/11_Spikes/raster_psth_all_good_units_dec4.png",
         take="Dec 4 only: dHPC + LEC at 50 Hz. Same readout: sorted single units, not LFP pickup; dHPC shows the strongest ON-window increase, while LEC is mixed and partly suppressive.",
         notes="This is the Dec 4-only all-unit view, with Dec 3 removed. Each row is one curated good unit at the 50 Hz / high-amplitude condition, averaged across 200 trials. Top is dorsal hippocampus, bottom is lateral entorhinal cortex. Red means firing above that unit's pre-ON baseline, black/gray means below. The right bars summarize ON-minus-OFF firing-rate change. The point is to show the full single-unit population without the Dec 3 control panel: dHPC carries the strongest positive ON-window response, while LEC is mixed and includes suppression. This is still the artifact-resistant signal because it is spike timing from sorted units, not the 50 Hz LFP line."),
    dict(kind="content", kicker="finding 6 · what the study shows", title="Region-specific, bidirectional modulation — active processing, not a passive relay",
         fig=f"{F}/11_Spikes/spike_50hz_interpretation.png",
         take="The study's bottom line: vibrotactile 50 Hz reaches these deep medial-temporal neurons and is processed region-specifically — hippocampus drives up, entorhinal cortex is suppressed. Opposite-signed responses to the same input rule out a passive relay, and pickup cannot remove spikes — so target engagement is established.",
         notes="This is the study's core conclusion. The two regions transform the same input in opposite directions: a hippocampal subset increases firing while entorhinal cortex is net-suppressed. Opposite-signed responses to identical input are inconsistent with a passive relay and consistent with circuit-specific processing — and the suppression is itself artifact-resistant (pickup adds spikes, it can't remove them). So the study establishes target engagement: a body-surface vibration at 50 Hz genuinely modulates deep medial-temporal neurons, region-specifically. The figure's inference ladder lays out the logic: not artifact, not a passive relay, but active region-specific processing."),
    # ===================== PART 4 — LIMITS & NEXT =====================
    dict(kind="section", title="Part 4 · Limits, and the next study"),
    dict(kind="content", kicker="the limitation", title="Limitations: entrainment was untestable — the sync channel never captured the carrier", fig=TTL,
         take="The study's main limitation: no time-aligned stimulus-phase reference was recorded, so phase-following / entrainment is untestable here — a measurement gap, not a null. The digital sync line updated at ~4 Hz independent of carrier (identical pulse counts at 5 vs 26 Hz; ~78 expected at 26 Hz), so no phase reference exists in the data. Separately, the lower-frequency stimuli were not clean sine waves, so 50-vs-lower specificity is partly confounded with stimulus quality.",
         notes="The limitations, together. Testing entrainment (phase-locking) requires a time-aligned reference for the stimulus phase, and we didn't record a usable one — so 'no entrainment shown' is a measurement gap, not evidence of absence. Verified from the raw digital stream: the sync channel updated at ~4 Hz irrespective of carrier (same ~6 pulses/trial for 5 and 26 Hz; ~78 expected at 26 Hz), never sustained a run at the carrier period — undersampled and decoupled, so phase is unrecoverable (a ~4 Hz observer can't resolve a 26 Hz oscillation). The power/periodicity negatives are real; only the phase-locking claim is gated by the missing reference. And post hoc we found the lower-frequency stimuli (5/10/26 Hz) were not clean sine waves, so the 50 Hz frequency-specificity is partly confounded with stimulus quality. The fix for the next study: record the stimulus redundantly on the acquisition clock (firmware per-cycle marker + transduced force/accelerometer + drive-signal copy), which makes entrainment — and eventually clearance readouts — testable."),

    dict(kind="closing", title="Thank you",
         sub="Happy to go into any figure, the artifact controls, the parameter plan, or the next-round instrumentation.",
         footer="wearables.snyder@gmail.com    ·    github.com/wehabit/hpaticmousebuzsakilab",
         notes="Anticipated questions. 'Is 50 Hz special or just strongest?' — strongest clean single-unit effect of the carriers tested; the LFP transient was largest at 26 Hz but that's the artifact-prone measure. 'Could it be arousal?' — possible, but frequency-specific (50 ≫ 26 at matched amplitude); an indirect sensory/state pathway isn't excluded. 'Why distrust the 50 Hz LFP?' — disconnected channels carry it at ~6× tissue. 'Is this a treatment?' — no; this is target-engagement evidence for a mechanism hypothesis, deliberately separated from any clinical claim."),
]


# ---------------------------------------------------------------- generated teaching figures
def build_evidence_ladder_figure(out: str | Path) -> None:
    """Visual glossary: claim-strength ladder (evoked -> steady-state -> phase locking -> entrainment).

    LFP vs single-unit is shown as how we READ the steady-state rung (a measurement
    choice), not as separate claim-strength steps: LFP clue -> Kilosort spike-sorting
    -> clean single-unit tuning.
    """
    import numpy as np
    import matplotlib
    matplotlib.use("Agg")
    import matplotlib.pyplot as plt
    from matplotlib.patches import Rectangle

    out = Path(out)
    out.parent.mkdir(parents=True, exist_ok=True)

    cardinal = TEAL
    navy = NAVY
    grey = GREY
    stone = STONE
    stone_dark = STONE_DARK
    stone_light = STONE_LIGHT

    def clean(ax):
        ax.set_facecolor("white")
        for spine in ax.spines.values():
            spine.set_visible(False)
        ax.set_xticks([]); ax.set_yticks([])
        return ax

    fig = plt.figure(figsize=(16, 9), dpi=220)
    fig.patch.set_facecolor("white")
    fig.text(0.04, 0.980, "Vocabulary first: from response to entrainment",
             fontsize=19, weight="bold", color=navy, va="top")
    fig.text(0.04, 0.940,
             "Four claim-strength steps. LFP vs single-unit is how we READ a step (a measurement choice), not a stronger claim.",
             fontsize=11.5, color=grey, va="top")

    rows = [
        ("1", "Evoked response",
         "One clap happens; the brain reacts.",
         "Data: signal changes at stimulation ON / OFF."),
        ("2", "Steady-state / frequency-following",
         "A steady beat appears as a peak at that rate. Read it left → right:",
         "LFP clue (pickup-prone) → sort spikes (Kilosort) → single-unit tuning (clean)."),
        ("3", "Phase locking",
         "Activity lands at the same beat-position every trial.",
         "Data: spikes / LFP align to the same delivered-stimulus phase (needs recorded phase)."),
        ("4", "Entrainment",
         "The outside beat captures an internal rhythm.",
         "Data: a native rhythm shifts into the external rhythm's timing."),
    ]

    y0 = 0.815
    row_h = 0.182
    chart_x = 0.66
    chart_w = 0.30
    fig.add_artist(plt.Line2D([0.045, 0.045],
                              [y0 + 0.045, y0 - (len(rows) - 1) * row_h - 0.055],
                              transform=fig.transFigure, color=stone_light, lw=2.0, alpha=0.75, zorder=0))

    for idx, (num, term, plain, data) in enumerate(rows):
        y = y0 - idx * row_h
        if idx > 0:
            fig.add_artist(plt.Line2D([0.075, 0.96], [y + 0.090, y + 0.090],
                                      transform=fig.transFigure, color=stone_light, lw=0.8, alpha=0.7, zorder=0))
        fig.text(0.045, y, num, fontsize=17, weight="bold", color="white", ha="center", va="center",
                 bbox=dict(boxstyle="circle,pad=0.30", facecolor=cardinal, edgecolor=cardinal))
        fig.text(0.085, y + 0.045, term, fontsize=14.5, weight="bold", color=navy, va="top")
        fig.text(0.085, y + 0.010, plain, fontsize=10.6, color=navy, va="top")
        fig.text(0.085, y - 0.022, data, fontsize=9.6, color=grey, va="top")

        if idx == 0:
            ax = clean(fig.add_axes([chart_x, y - 0.045, chart_w, 0.095]))
            t = np.linspace(-1, 4, 600)
            sig = 0.05 * np.sin(2 * np.pi * 1.5 * t)
            sig += 0.95 * np.exp(-((t - 0.05) / 0.15) ** 2)
            sig += 0.55 * np.exp(-((t - 3.0) / 0.18) ** 2)
            ax.axvspan(0, 3, color=cardinal, alpha=0.11)
            ax.axvline(0, color=cardinal, lw=1.5); ax.axvline(3, color=cardinal, lw=1.5)
            ax.plot(t, sig, color=navy, lw=2)
            ax.text(0.03, 0.84, "ON", transform=ax.transAxes, fontsize=8, color=cardinal)
            ax.text(0.72, 0.84, "OFF", transform=ax.transAxes, fontsize=8, color=cardinal)
            ax.set_xlim(-0.7, 3.7); ax.set_ylim(-0.35, 1.2)
        elif idx == 1:
            # three reads of steady-state, left -> right: LFP line | Kilosort | single-unit tuning
            axL = clean(fig.add_axes([0.600, y - 0.046, 0.100, 0.092]))
            freq = np.linspace(5, 95, 600)
            background = 1.25 / np.sqrt(freq)
            power = background + 0.48 * np.exp(-0.5 * ((freq - 50) / 2.8) ** 2)
            axL.plot(freq, background, color="#B8B8B8", lw=1.2, ls="--")
            axL.plot(freq, power, color=stone, lw=1.8)
            axL.axvline(50, color=cardinal, lw=1.3)
            axL.fill_between(freq, background, power, where=(freq > 45) & (freq < 55), color=cardinal, alpha=0.18)
            axL.text(0.5, 1.16, "LFP: 50 Hz line", transform=axL.transAxes, fontsize=7.3, color=stone_dark, ha="center", weight="bold")
            axL.text(0.5, 0.97, "(pickup-prone)", transform=axL.transAxes, fontsize=6.6, color=grey, ha="center")
            axL.set_xlim(5, 95); axL.set_ylim(0, 0.86)

            fig.text(0.713, y, "→", fontsize=15, color=grey, ha="center", va="center")

            axM = clean(fig.add_axes([0.728, y - 0.046, 0.108, 0.092]))
            axM.set_xlim(0, 1); axM.set_ylim(0, 1)
            offsets = np.array([[-0.05, -0.02], [-0.03, 0.05], [0.0, 0.0],
                                [0.04, -0.04], [0.05, 0.03], [-0.01, -0.06]]) * 1.5
            for cx, cy, col in [(0.30, 0.60, stone_dark), (0.52, 0.36, cardinal), (0.70, 0.66, stone)]:
                pts = offsets + np.array([cx, cy])
                axM.scatter(pts[:, 0], pts[:, 1], s=11, color=col, alpha=0.82, edgecolors="none")
            axM.text(0.5, 1.16, "Kilosort:", transform=axM.transAxes, fontsize=7.3, color=stone_dark, ha="center", weight="bold")
            axM.text(0.5, 0.97, "sort spikes", transform=axM.transAxes, fontsize=6.6, color=grey, ha="center")

            fig.text(0.856, y, "→", fontsize=15, color=grey, ha="center", va="center")

            axR = clean(fig.add_axes([0.872, y - 0.044, 0.104, 0.088]))
            vals = np.array([[0.05, 0.03, 0.18, 0.96], [0.04, 0.02, 0.12, 0.55], [0.06, 0.02, 0.08, 0.24]])
            axR.imshow(vals, cmap="Reds", vmin=0, vmax=1, aspect="auto")
            axR.set_xticks(range(4)); axR.set_xticklabels(["5", "10", "26", "50"], fontsize=6.6)
            axR.set_yticks(range(3)); axR.set_yticklabels(["amp250", "amp180", "amp100"], fontsize=6.2)
            axR.tick_params(length=0, pad=1)
            for c in range(4):
                for r in range(3):
                    axR.add_patch(plt.Rectangle((c - 0.5, r - 0.5), 1, 1, fill=False, edgecolor="white", lw=0.6))
            axR.text(0.5, 1.22, "single units: freq x amp", transform=axR.transAxes, fontsize=7.0, color=stone_dark, ha="center", weight="bold")
            axR.text(3, 0, "biggest", fontsize=6.6, color="white", ha="center", va="center", weight="bold")
        elif idx == 2:
            ax = clean(fig.add_axes([chart_x, y - 0.052, chart_w, 0.11]))
            lock_times = np.array([0.08, 0.42, 0.76])
            tt = np.linspace(0, 1, 800)
            ax.plot(tt, 0.24 * np.sin(2 * np.pi * 3 * tt) + 0.80, color="#4F565A", lw=1.6)
            for p in lock_times:
                ax.plot([p], [1.04], marker="v", color=stone_dark, ms=5, zorder=3)
            for trial, ydot in enumerate([0.34, 0.16, -0.02, -0.20]):
                ax.hlines(ydot, 0.0, 1.0, color="#D9D9D9", lw=1.0, zorder=0)
                for p, j in zip(lock_times, [-0.004, 0.003, -0.002]):
                    ax.vlines(p + j + trial * 0.001, ydot - 0.06, ydot + 0.06, color=stone_dark, lw=3.2, zorder=4)
                ax.text(-0.03, ydot, f"trial {trial + 1}", fontsize=6.8, color=navy, ha="right", va="center")
            ax.text(0.5, -0.34, "spike ticks line up = same phase", transform=ax.transAxes,
                    fontsize=7.6, color=stone_dark, ha="center", va="top", weight="bold")
            ax.set_xlim(0, 1); ax.set_ylim(-0.42, 1.14)
        else:
            ax = clean(fig.add_axes([chart_x, y - 0.055, chart_w, 0.12]))
            tt = np.linspace(0, 1, 800); stim_on = 0.38
            external = 0.34 * np.sin(2 * np.pi * 4.5 * (tt - stim_on))
            native = np.where(tt < stim_on, 0.18 * np.sin(2 * np.pi * 2.4 * tt + 1.1),
                              0.58 * np.sin(2 * np.pi * 4.5 * (tt - stim_on) - 0.05))
            ax.axvspan(stim_on, 1.0, color=cardinal, alpha=0.08)
            ax.plot(tt[tt >= stim_on], external[tt >= stim_on] + 0.62, color=cardinal, lw=1.4)
            ax.plot(tt, native - 0.30, color=stone_dark, lw=2.0)
            ax.axvline(stim_on, color=navy, lw=1, alpha=0.45)
            ax.text(0.02, 0.18, "own pace", transform=ax.transAxes, fontsize=7.6, color=stone_dark)
            ax.text(0.52, 0.16, "captured", transform=ax.transAxes, fontsize=7.6, color=stone_dark)
            ax.text(0.5, 0.86, "external beat ON", transform=ax.transAxes, fontsize=7.6, color=cardinal, ha="center")
            ax.set_xlim(0, 1); ax.set_ylim(-1.08, 1.18)

    fig.add_artist(Rectangle((0.03, 0.030), 0.94, 0.062, transform=fig.transFigure,
                             facecolor=FOG_LIGHT, edgecolor=stone_light, linewidth=1.0))
    fig.text(0.05, 0.068,
             "Use this vocabulary on the next slide: a stimulus-locked line can be target engagement without proving native entrainment.",
             fontsize=11.7, weight="bold", color=navy, va="center")
    fig.text(0.05, 0.042,
             "Our haptic data are read on this same ladder: response first; phase-locking and entrainment require stronger timing evidence.",
             fontsize=10.1, color=grey, va="center")
    fig.savefig(out, bbox_inches="tight", pad_inches=0.08)
    plt.close(fig)


# ---------------------------------------------------------------- PPTX
def build_pptx(slides, out):
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    from pptx.enum.shapes import MSO_SHAPE
    from PIL import Image

    def rgb(h):
        return RGBColor(int(h[1:3], 16), int(h[3:5], 16), int(h[5:7], 16))
    NV, TL, GR, PN, WH = rgb(NAVY), rgb(TEAL), rgb(GREY), rgb(PANEL), RGBColor(255, 255, 255)
    SW, SH = 13.333, 7.5
    prs = Presentation(); prs.slide_width = Inches(SW); prs.slide_height = Inches(SH)
    BL = prs.slide_layouts[6]

    def txt(s, x, y, w, h, text, size, *, bold=False, color=NV, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
        tb = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
        tf = tb.text_frame; tf.word_wrap = True; tf.vertical_anchor = anchor
        for i, line in enumerate(text.split("\n")):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = line; p.alignment = align
            r = p.runs[0]; r.font.size = Pt(size); r.font.bold = bold; r.font.color.rgb = color; r.font.name = "Calibri"

    def bar(s, x, y, w, h, color):
        sh = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
        sh.fill.solid(); sh.fill.fore_color.rgb = color; sh.line.fill.background()

    def img(s, path, x, y, width):
        if Path(path).exists():
            w, h = Image.open(path).size
            s.shapes.add_picture(path, Inches(x), Inches(y), width=Inches(width), height=Inches(width * h / w))

    def pic(s, path, maxw, maxh, top):
        if not Path(path).exists():
            txt(s, 1, top, SW - 2, 1, f"[missing: {path}]", 14, color=GR); return
        w, h = Image.open(path).size; ar = w / h
        iw, ih = maxw, maxw / ar
        if ih > maxh:
            ih, iw = maxh, maxh * ar
        s.shapes.add_picture(str(path), Inches((SW - iw) / 2), Inches(top), width=Inches(iw), height=Inches(ih))

    def pic_box(s, path, x, y, maxw, maxh):
        if not Path(path).exists():
            txt(s, x, y + maxh / 2 - 0.25, maxw, 0.5, f"[missing: {path}]", 12, color=GR, align=PP_ALIGN.CENTER)
            return
        w, h = Image.open(path).size; ar = w / h
        iw, ih = maxw, maxw / ar
        if ih > maxh:
            ih, iw = maxh, maxh * ar
        s.shapes.add_picture(str(path), Inches(x + (maxw - iw) / 2), Inches(y + (maxh - ih) / 2),
                             width=Inches(iw), height=Inches(ih))

    for i, sl in enumerate(slides):
        s = prs.slides.add_slide(BL); k = sl["kind"]
        if k in ("title", "closing"):
            bar(s, 0, 0, SW, SH, NV); bar(s, 0, 5.0, SW, 0.08, TL); img(s, LOGO_W, 0.55, 0.5, 3.6)
            if k == "title":
                img(s, MOUSE, SW - 2.9, SH - 1.7, 2.2)
            txt(s, 1, 2.0 if k == "title" else 3.0, SW - 4.0, 1.6, sl["title"], 34, bold=True, color=WH, anchor=MSO_ANCHOR.MIDDLE)
            if sl.get("sub"):
                txt(s, 1, 3.7 if k == "title" else 4.2, SW - 4.0, 1.2, sl["sub"], 17, color=rgb("#D7D2CB"))
            if sl.get("sub2"):
                txt(s, 1, 4.25, SW - 4.0, 0.6, sl["sub2"], 16, color=rgb("#9A958C"))
            if sl.get("footer"):
                txt(s, 1, 5.25, SW - 4.0, 0.35, sl["footer"], 15.5, color=rgb("#9A958C"))
            if sl.get("footer2"):
                txt(s, 1, 5.68, SW - 4.0, 0.35, sl["footer2"], 15.5, color=rgb("#D7D2CB"))
        elif k == "section":
            bar(s, 0, 0, SW, SH, TL); img(s, LOGO_W, 0.55, 0.5, 2.9)
            txt(s, 1, 2.7, SW - 2, 1.5, sl["title"], 34, bold=True, color=WH, anchor=MSO_ANCHOR.MIDDLE)
            if sl.get("sub"):
                sub_y = sl.get("sub_y", 4.1)
                txt(s, 1, sub_y, SW - 2, 0.65, sl["sub"], 17, color=rgb("#EDDDDD"))
            if sl.get("refs"):
                txt(s, 1, sl.get("refs_y", sl.get("sub_y", 4.1) + 0.55), SW - 2, 0.5,
                    sl["refs"], 10.5, color=rgb("#D7D2CB"))
        elif k == "side_by_side":
            bar(s, 0, 0, SW, 0.18, TL); y = 0.45
            if sl.get("kicker"):
                txt(s, 0.6, y, SW - 1.2, 0.4, sl["kicker"].upper(), 14, bold=True, color=TL); y += 0.45
            txt(s, 0.6, y, SW - 1.2, 0.72, sl["title"], 23, bold=True, color=NV)
            left_x, right_x = 0.45, 6.85
            panel_y, panel_w, panel_h = 1.95, 5.95, 4.42
            for x, (path, label) in zip((left_x, right_x), sl["figs"]):
                txt(s, x, 1.55, panel_w, 0.32, label, 12.5, bold=True, color=TL, align=PP_ALIGN.CENTER)
                pic_box(s, path, x, panel_y, panel_w, panel_h)
            if sl.get("take"):
                bar(s, 0, SH - 0.9, SW, 0.9, PN)
                txt(s, 0.6, SH - 0.87, SW - 1.2, 0.85, sl["take"], 14.5, bold=True, color=TL, anchor=MSO_ANCHOR.MIDDLE)
        elif sl.get("full_fig"):
            if Path(sl["fig"]).exists():
                s.shapes.add_picture(str(sl["fig"]), Inches(0.15), Inches(0.10), width=Inches(SW - 0.30))
            else:
                txt(s, 1, 2, SW - 2, 1, f"[missing: {sl['fig']}]", 16, color=GR)
        else:
            bar(s, 0, 0, SW, 0.18, TL); y = 0.45
            if sl.get("kicker"):
                txt(s, 0.6, y, SW - 1.2, 0.4, sl["kicker"].upper(), 14, bold=True, color=TL); y += 0.45
            txt(s, 0.6, y, SW - 1.2, 1.1, sl["title"], 24, bold=True, color=NV)
            if sl.get("video"):
                # split layout: portrait setup video on the LEFT, figure on the RIGHT
                top = y + 1.25
                ch = (SH - 0.9) - top - 0.05            # content height above the take bar
                vcol = 2.7                              # left column reserved for the video
                fx = 0.6 + vcol + 0.45                  # figure starts right of the video column
                pic_box(s, sl["fig"], fx, top, SW - fx - 0.5, ch)
                vw, vh = sl.get("video_size", (1080, 1920))
                ih = ch; iw = ih * vw / vh
                if iw > vcol:
                    iw = vcol; ih = iw * vh / vw
                vx = 0.6 + (vcol - iw) / 2
                vy = top + (ch - ih) / 2
                if sl.get("video_embed") and Path(sl["video"]).exists():
                    # embed the playable movie (heavy; can choke the LibreOffice PDF render)
                    s.shapes.add_movie(sl["video"], Inches(vx), Inches(vy), Inches(iw), Inches(ih),
                                       poster_frame_image=sl.get("video_poster"),
                                       mime_type="video/quicktime")
                elif sl.get("video_poster") and Path(sl["video_poster"]).exists():
                    # placeholder: show the setup still; drop the real video onto it in PowerPoint
                    pic_box(s, sl["video_poster"], vx, vy, iw, ih)
                    txt(s, vx + iw / 2, vy - 0.16, iw + 1.0, 0.3, "▶ setup video — insert in PowerPoint",
                        9, bold=True, color=TL, align=PP_ALIGN.CENTER)
                else:
                    txt(s, vx, vy + ih / 2, iw, 0.5, "[video missing]", 12, color=GR, align=PP_ALIGN.CENTER)
            elif sl.get("fig"):
                import concept_native
                stem = Path(sl["fig"]).stem
                if not sl.get("full_fig") and stem in concept_native.NATIVE:
                    by = y + 1.25
                    concept_native.render(s, stem, 0.8, by, SW - 1.6, (SH - 0.9) - by - 0.05)
                elif sl.get("full_fig"):
                    pic(s, sl["fig"], SW - 0.7, SH - (y + 0.95) - 0.28, top=y + 0.75)
                else:
                    pic(s, sl["fig"], sl.get("fig_maxw", SW - 1.6),
                        SH - (y + 1.4) - 0.9, top=y + 1.25)
            if sl.get("take"):
                bar(s, 0, SH - 0.9, SW, 0.9, PN)
                txt(s, 0.6, SH - 0.87, SW - 1.2, 0.85, sl["take"], 15, bold=True, color=TL, anchor=MSO_ANCHOR.MIDDLE)
        if sl.get("notes"):
            s.notes_slide.notes_text_frame.text = sl["notes"]
        if k != "title":
            txt(s, SW - 1.8, 0.27, 1.4, 0.4, f"{i + 1} / {len(slides)}", 11,
                color=(WH if k in ("section", "closing") else GR), align=PP_ALIGN.RIGHT)
    prs.save(out)


# ---------------------------------------------------------------- PDF
def build_pdf(slides, out):
    import matplotlib
    matplotlib.use("Agg")
    import matplotlib.pyplot as plt
    import matplotlib.image as mpimg
    from PIL import Image
    from tempfile import TemporaryDirectory
    slide_h = 7.5
    # Render full slides to high-resolution page images before assembling the PDF.
    # Matplotlib's native PDF backend can look soft in some viewers for this deck's
    # dense raster figures, even when the source PNGs are high resolution.
    pdf_dpi = 300

    def wrap(t, w):
        return "\n".join(textwrap.fill(line, w) for line in t.split("\n"))

    def place(fig, path, box):
        if not Path(path).exists():
            fig.text(0.5, box[1] + box[3] / 2, f"[missing: {path}]", ha="center", color=GREY); return
        im = mpimg.imread(path); ar = im.shape[1] / im.shape[0]
        fw, fh = fig.get_size_inches()
        bw_in, bh_in = box[2] * fw, box[3] * fh
        dw_in = min(bw_in, bh_in * ar); dh_in = dw_in / ar
        wf, hf = dw_in / fw, dh_in / fh
        ax = fig.add_axes([box[0] + (box[2] - wf) / 2, box[1] + (box[3] - hf) / 2, wf, hf])
        ax.imshow(im, interpolation="none"); ax.axis("off")

    def imgbox(fig, path, x, top, w):
        if not Path(path).exists():
            return
        im = mpimg.imread(path); ar = im.shape[1] / im.shape[0]
        fw, fh = fig.get_size_inches(); hf = (w * fw / ar) / fh
        ax = fig.add_axes([x, top - hf, w, hf]); ax.imshow(im, interpolation="none"); ax.axis("off")

    page_paths = []
    with TemporaryDirectory() as tmpdir:
        for i, sl in enumerate(slides):
            fig = plt.figure(figsize=(13.333, 7.5), dpi=pdf_dpi); k = sl["kind"]
            if k in ("title", "closing"):
                fig.patch.set_facecolor(NAVY); imgbox(fig, LOGO_W, 0.05, 0.93, 0.30)
                if k == "title":
                    imgbox(fig, MOUSE, 0.80, 0.22, 0.16)
                fig.text(0.07, 0.58, wrap(sl["title"], 30), fontsize=32, color="white", weight="bold", va="center")
                if sl.get("sub"):
                    fig.text(0.07, 0.36, wrap(sl["sub"], 74), fontsize=15, color="#D7D2CB")
                if sl.get("sub2"):
                    fig.text(0.07, 0.295, wrap(sl["sub2"], 80), fontsize=14, color="#9A958C")
                if sl.get("footer"):
                    fig.text(0.07, 0.16, sl["footer"], fontsize=13, color="#9A958C")
                if sl.get("footer2"):
                    fig.text(0.07, 0.105, sl["footer2"], fontsize=13, color="#D7D2CB")
            elif k == "section":
                fig.patch.set_facecolor(TEAL); imgbox(fig, LOGO_W, 0.05, 0.93, 0.24)
                fig.text(0.08, 0.56, wrap(sl["title"], 34), fontsize=33, color="white", weight="bold", va="center")
                if sl.get("sub"):
                    sub_y = 0.40 if "sub_y" not in sl else max(0.13, 1.0 - (sl["sub_y"] / slide_h))
                    fig.text(0.08, sub_y, wrap(sl["sub"], 78), fontsize=15, color="#EDDDDD")
                if sl.get("refs"):
                    refs_y = max(0.08, 1.0 - (sl.get("refs_y", sl.get("sub_y", 4.1) + 0.55) / slide_h))
                    fig.text(0.08, refs_y, wrap(sl["refs"], 102), fontsize=9.2, color="#D7D2CB")
            elif k == "side_by_side":
                fig.patch.set_facecolor("white")
                fig.add_artist(plt.Rectangle((0, 0.975), 1, 0.025, color=TEAL, transform=fig.transFigure))
                yt = 0.93
                if sl.get("kicker"):
                    fig.text(0.05, 0.945, sl["kicker"].upper(), fontsize=12, color=TEAL, weight="bold"); yt = 0.90
                fig.text(0.05, yt, wrap(sl["title"], 70), fontsize=21, color=NAVY, weight="bold", va="top")
                boxes = [(0.035, 0.17, 0.445, 0.62), (0.520, 0.17, 0.445, 0.62)]
                for (path, label), box in zip(sl["figs"], boxes):
                    fig.text(box[0] + box[2] / 2, 0.805, label, fontsize=11.5, color=TEAL, weight="bold", ha="center")
                    place(fig, path, box)
                if sl.get("take"):
                    fig.add_artist(plt.Rectangle((0, 0), 1, 0.12, color=PANEL, transform=fig.transFigure))
                    fig.text(0.05, 0.06, wrap(sl["take"], 120), fontsize=12.5, color=TEAL, weight="bold", va="center")
            elif sl.get("full_fig"):
                fig.patch.set_facecolor("white")
                place(fig, sl["fig"], (0.01, 0.01, 0.98, 0.98))
            else:
                fig.patch.set_facecolor("white")
                fig.add_artist(plt.Rectangle((0, 0.975), 1, 0.025, color=TEAL, transform=fig.transFigure))
                yt = 0.93
                if sl.get("kicker"):
                    fig.text(0.05, 0.945, sl["kicker"].upper(), fontsize=12, color=TEAL, weight="bold"); yt = 0.90
                fig.text(0.05, yt, wrap(sl["title"], 70), fontsize=21, color=NAVY, weight="bold", va="top")
                if sl.get("video"):
                    # match build_pptx: portrait setup-video poster on the LEFT, figure on the RIGHT
                    if sl.get("video_poster") and Path(sl["video_poster"]).exists():
                        place(fig, sl["video_poster"], (0.045, 0.15, 0.185, 0.62))
                        fig.text(0.1375, 0.80, "▶ setup video — insert in PowerPoint",
                                 fontsize=7.4, color=TEAL, weight="bold", ha="center")
                    place(fig, sl["fig"], (0.255, 0.15, 0.715, 0.63))
                elif sl.get("fig"):
                    if sl.get("full_fig"):
                        place(fig, sl["fig"], (0.025, 0.035, 0.95, 0.78))
                    else:
                        place(fig, sl["fig"], sl.get("fig_box", (0.05, 0.18, 0.90, 0.60)))
                if sl.get("take"):
                    fig.add_artist(plt.Rectangle((0, 0), 1, 0.12, color=PANEL, transform=fig.transFigure))
                    fig.text(0.05, 0.06, wrap(sl["take"], 120), fontsize=13, color=TEAL, weight="bold", va="center")
            if k != "title":
                fig.text(0.955, 0.955, f"{i + 1} / {len(slides)}", ha="right", va="top", fontsize=11,
                         color=("white" if k in ("section", "closing") else GREY))
            page_path = Path(tmpdir) / f"slide_{i + 1:02d}.jpg"
            fig.savefig(page_path, format="jpg", facecolor=fig.get_facecolor(), dpi=pdf_dpi,
                        pil_kwargs={"quality": 96, "subsampling": 0})
            plt.close(fig)
            page_paths.append(page_path)

        page_images = [Image.open(p).convert("RGB") for p in page_paths]
        try:
            page_images[0].save(out, save_all=True, append_images=page_images[1:],
                                resolution=pdf_dpi, quality=96, subsampling=0)
        finally:
            for im in page_images:
                im.close()


def main():
    pptx = OUT / "haptic_brain_talk.pptx"; pdf = OUT / "haptic_brain_talk.pdf"
    build_evidence_ladder_figure(EVIDENCE)
    build_pptx(SLIDES, pptx); print("wrote", pptx)
    build_pdf(SLIDES, pdf); print("wrote", pdf)
    print(f"{len(SLIDES)} slides")


if __name__ == "__main__":
    main()

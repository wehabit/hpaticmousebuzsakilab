"""Native-PowerPoint renderers for the concept/teaching diagrams.

Instead of pasting a matplotlib PNG/SVG, these draw the diagrams as REAL PowerPoint
objects (rounded rectangles, text boxes, arrows) so every element is directly editable
in PowerPoint — no "convert to shape", infinitely crisp. Coordinates are translated from
each figure's original matplotlib data space into the slide box it occupies.

Only the box-and-arrow diagrams live here; data plots (heatmaps, spectrograms) stay as
high-resolution images because they are plotted data, not shapes.
"""
from __future__ import annotations

import math

from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn

# Stanford palette (mirrors build_concept_figures)
BLACK = "#2E2D29"; CARDINAL = "#8C1515"; COOL_GREY = "#53565A"; FOG = "#DAD7CB"
FOG_LIGHT = "#F4F4F4"; STONE = "#7F7776"; STONE_DARK = "#544948"; LAGUNITA = "#007C92"
POPPY = "#E98300"; PALO_ALTO = "#175E54"
NAVY, TEAL, GOLD, RED, GREEN, GREY, LIGHT = BLACK, STONE, STONE_DARK, CARDINAL, STONE_DARK, COOL_GREY, FOG


def _rgb(h):
    return RGBColor(int(h[1:3], 16), int(h[3:5], 16), int(h[5:7], 16))


class Pen:
    """Maps a figure's data coords (0..XMAX, 0..YMAX, y-up) into a slide box (inches)."""

    def __init__(self, s, xmax, ymax, bx, by, bw, bh):
        self.s = s; self.xmax = xmax; self.ymax = ymax
        # fit preserving aspect ratio, centered in the available box
        ar = xmax / ymax
        if bw / bh > ar:
            dw, dh = bh * ar, bh
        else:
            dw, dh = bw, bw / ar
        self.x0 = bx + (bw - dw) / 2; self.y0 = by + (bh - dh) / 2
        self.sx = dw / xmax; self.sy = dh / ymax

    def X(self, x): return self.x0 + x * self.sx
    def Y(self, y): return self.y0 + (self.ymax - y) * self.sy   # flip y
    def W(self, w): return w * self.sx
    def H(self, h): return h * self.sy
    def pt(self, mpl_fs): return max(6.0, mpl_fs * self.sx)      # mpl figs were ~1 in / data unit

    def box(self, x, y, w, h, fc=None, ec=None, lw=1.5, rounded=True):
        shp = MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE
        b = self.s.shapes.add_shape(shp, Inches(self.X(x)), Inches(self.Y(y + h)),
                                    Inches(self.W(w)), Inches(self.H(h)))
        if fc is None:
            b.fill.background()
        else:
            b.fill.solid(); b.fill.fore_color.rgb = _rgb(fc)
        if ec is None:
            b.line.fill.background()
        else:
            b.line.color.rgb = _rgb(ec); b.line.width = Pt(lw)
        b.shadow.inherit = False
        b.text_frame.word_wrap = True
        return b

    def text(self, x, y, txt, fs, color=NAVY, bold=False, align="center", anchor="middle",
             italic=False, w=3.0, h=0.7):
        wI, hI = self.W(w), self.H(h)               # w, h given in DATA units
        left = self.X(x) - wI / 2 if align == "center" else self.X(x)
        tb = self.s.shapes.add_textbox(Inches(left), Inches(self.Y(y) - hI / 2), Inches(wI), Inches(hI))
        tf = tb.text_frame; tf.word_wrap = True
        tf.vertical_anchor = {"middle": MSO_ANCHOR.MIDDLE, "top": MSO_ANCHOR.TOP, "bottom": MSO_ANCHOR.BOTTOM}[anchor]
        tf.margin_left = 0; tf.margin_right = 0; tf.margin_top = 0; tf.margin_bottom = 0
        al = {"center": PP_ALIGN.CENTER, "left": PP_ALIGN.LEFT, "right": PP_ALIGN.RIGHT}[align]
        for i, line in enumerate(txt.split("\n")):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.alignment = al; r = p.add_run(); r.text = line
            r.font.size = Pt(self.pt(fs)); r.font.bold = bold; r.font.italic = italic
            r.font.color.rgb = _rgb(color); r.font.name = "Calibri"
        return tb

    def textblock(self, x, ytop, lines, fs, color, w=3.0, h=1.6):
        """Multi-line block (each line: (text, bold)) in one top-anchored textbox."""
        wI, hI = self.W(w), self.H(h)
        tb = self.s.shapes.add_textbox(Inches(self.X(x) - wI / 2), Inches(self.Y(ytop)), Inches(wI), Inches(hI))
        tf = tb.text_frame; tf.word_wrap = True; tf.vertical_anchor = MSO_ANCHOR.TOP
        tf.margin_left = 0; tf.margin_right = 0; tf.margin_top = 0; tf.margin_bottom = 0
        for i, (line, bold) in enumerate(lines):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.alignment = PP_ALIGN.CENTER; r = p.add_run(); r.text = line
            r.font.size = Pt(self.pt(fs)); r.font.bold = bold
            r.font.color.rgb = _rgb(color); r.font.name = "Calibri"
        return tb

    def arrow(self, x0, y0, x1, y1, color=NAVY, lw=2.4):
        c = self.s.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(self.X(x0)), Inches(self.Y(y0)),
                                        Inches(self.X(x1)), Inches(self.Y(y1)))
        c.line.color.rgb = _rgb(color); c.line.width = Pt(lw)
        ln = c.line._get_or_add_ln()
        te = ln.makeelement(qn("a:tailEnd"), {"type": "triangle", "w": "med", "len": "med"})
        ln.append(te)
        return c

    def line(self, x0, y0, x1, y1, color=NAVY, lw=2.0, dash=None):
        c = self.s.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(self.X(x0)), Inches(self.Y(y0)),
                                        Inches(self.X(x1)), Inches(self.Y(y1)))
        c.line.color.rgb = _rgb(color); c.line.width = Pt(lw)
        if dash:
            ln = c.line._get_or_add_ln()
            ln.append(ln.makeelement(qn("a:prstDash"), {"val": dash}))
        return c

    def oval(self, cx, cy, rw, rh, fc=None, ec=None, lw=1.5):
        b = self.s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(self.X(cx - rw)), Inches(self.Y(cy + rh)),
                                    Inches(self.W(2 * rw)), Inches(self.H(2 * rh)))
        if fc is None:
            b.fill.background()
        else:
            b.fill.solid(); b.fill.fore_color.rgb = _rgb(fc)
        if ec is None:
            b.line.fill.background()
        else:
            b.line.color.rgb = _rgb(ec); b.line.width = Pt(lw)
        b.shadow.inherit = False
        return b

    def curve(self, pts, color=NAVY, lw=2.0):   # polyline via straight segments (robust)
        for (a, b) in zip(pts[:-1], pts[1:]):
            self.line(a[0], a[1], b[0], b[1], color=color, lw=lw)


# ---------------------------------------------------------------- figure renderers
def whitespace(p: Pen):
    p.text(5.5, 3.95, "GENUS clinical context and our wearable haptic path", 16, NAVY, bold=True, w=10)
    p.box(0.35, 0.30, 4.95, 3.00, fc=LIGHT, ec=GREY, lw=1.5)
    p.text(2.825, 3.00, "Audiovisual 40 Hz (GENUS)", 13.5, GREY, bold=True, w=4.7)
    p.text(2.825, 2.62, "Cognito Spectris / GammaSense device", 9.3, GREY, bold=True, w=4.65)
    p.text(2.825, 2.40, "Alzheimer's disease trials (mostly mild-to-moderate AD)", 6.8, GREY, w=4.65)
    p.text(0.70, 2.05,
           "OVERTURE (n=76): safe/tolerable;\nprimary endpoint did not separate;\nexploratory ADCS-ADL/MMSE/MRI signals",
           5.8, GREY, align="left", w=3.38, h=0.62)
    p.text(4.35, 2.05, "feasibility\npublished", 5.6, RED, bold=True, w=0.95, h=0.42)
    p.text(0.70, 1.38,
           "HOPE (n=670): testing 12-mo ADCS-ADL/MMSE slowing\nOLE (n=402): all-active extension",
           5.8, GREY, align="left", w=3.38, h=0.48)
    p.text(4.35, 1.38, "pivotal +\nextension", 5.6, RED, bold=True, w=0.95, h=0.42)
    p.text(0.70, 0.93,
           "ETUDE (n=20): dosing/safety/amyloid PET;\nno posted results",
           5.8, GREY, align="left", w=3.38, h=0.42)
    p.text(4.35, 0.93, "dose-\nranging", 5.6, RED, bold=True, w=0.95, h=0.42)
    p.text(0.70, 0.52,
           "Mlinarič 2025: human iEEG target engagement\n(n=11 epilepsy; 490 contacts; independent)",
           5.7, GREY, align="left", w=3.95, h=0.36)
    p.text(4.35, 0.52, "hippocampus\nengagement", 5.6, RED, bold=True, w=0.95, h=0.36)
    p.box(6.1, 0.30, 4.4, 3.00, fc=LIGHT, ec=RED, lw=2.6)
    p.text(8.3, 3.00, "Haptics (us)", 14, RED, bold=True, w=4.2)
    p.text(8.3, 2.62, "Vibrotactile device applied to fingertips", 9.3, NAVY, bold=True, w=4.05)
    p.text(8.3, 2.12,
           "Benefits: wearable for long sessions;\nsleep-compatible; pairs with sensing",
           7.7, RED, bold=True, w=4.05, h=0.5)
    p.text(8.3, 1.52, "2-finger glove:\nindex + middle finger", 8.9, NAVY, w=4.05, h=0.5)
    p.text(8.3, 0.80,
           "Potential human testing:\nStanford Hospital iEEG, stroke, or AD patients\n(resource/funding-driven)",
           6.6, NAVY, w=4.1, h=0.7)


def vision_chain(p: Pen):
    p.text(6, 3.15, "The proposal — a mechanism hypothesis to test", 15, NAVY, bold=True, w=11)
    boxes = [("Wearable\nvibrotactile\nstimulation", RED), ("Brain rhythms\n(entrainment)", NAVY),
             ("Glymphatic\nclearance ↑", NAVY), ("Amyloid-β /\ntau cleared", RED)]
    w, gap, y, h = 2.4, 0.7, 1.05, 1.5
    xs = [0.15 + k * (w + gap) for k in range(4)]
    for (lab, tcol), x in zip(boxes, xs):
        p.box(x, y, w, h, fc=FOG_LIGHT, ec=tcol, lw=2.2)
        p.text(x + w / 2, y + h / 2, lab, 12.5, tcol, bold=True, w=w)
    for k in range(3):
        p.arrow(xs[k] + w, y + h / 2, xs[k + 1], y + h / 2, color=GREY, lw=2.6)
    p.text(6, 0.45, "a mechanism hypothesis — to be tested, not yet a treatment claim", 11, GREY, italic=True, w=11)


def genus_evidence_chain(p: Pen):
    p.text(5.75, 3.95, "What each 40 Hz anchor adds", 15, NAVY, bold=True, w=11)
    p.text(5.75, 3.65, "mouse target/pathology -> glymphatic mechanism -> synthesis -> early human extension",
           8.8, GREY, w=11, h=0.3)

    cards = [
        ("Tsai / GENUS\n2016-2019",
         "40 Hz light/sound\ndrives gamma activity;\nreduces Aβ/tau pathology;\nengages microglia + circuits",
         "mouse pathology"),
        ("Murdock\nNature 2024",
         "40 Hz audiovisual\nincreases CSF influx\nand ISF efflux;\nAQP4-dependent\namyloid clearance",
         "glymphatic bridge"),
        ("Park & Tsai\nPLOS Biol 2025",
         "Review synthesizes\npreclinical mechanisms\nand early clinical evidence;\nflags open questions",
         "field synthesis"),
        ("Chan et al.\nAlz Dement 2025",
         "Small mild-AD\nopen-label extension:\nfeasible/safe;\nsuggestive longer-term\nsignals, no blinded control",
         "early human"),
    ]
    w, h, gap, y = 2.45, 2.25, 0.28, 1.05
    xs = [0.35 + i * (w + gap) for i in range(4)]
    for i, ((head, body, tag), x) in enumerate(zip(cards, xs), start=1):
        edge = RED if i in (1, 2) else GREY
        p.box(x, y, w, h, fc=FOG_LIGHT, ec=edge, lw=2.0)
        p.text(x + 0.20, y + h - 0.22, str(i), 8.6, RED, bold=True, align="left", w=0.35, h=0.28)
        p.text(x + w / 2, y + h - 0.46, head, 9.0, RED if i == 2 else NAVY, bold=True, w=w - 0.25, h=0.48)
        p.text(x + w / 2, y + 1.13, body, 6.8, STONE_DARK, bold=True, w=w - 0.25, h=1.00)
        p.text(x + w / 2, y + 0.24, tag, 7.2, RED, bold=True, w=w - 0.25, h=0.28)
    for i in range(3):
        p.arrow(xs[i] + w + 0.05, y + h / 2, xs[i + 1] - 0.05, y + h / 2, color=GREY, lw=2.1)
    p.text(5.75, 0.38,
           "takeaway: strong mouse mechanism; human efficacy remains an active controlled-trial question",
           9.2, RED, bold=True, w=11, h=0.35)


def response_entrainment_guardrail(p: Pen):
    p.text(5.75, 3.95, "A 40 Hz response can be real without proving native gamma entrainment",
           13.8, NAVY, bold=True, w=11)
    p.text(5.75, 3.64, "Buzsáki/Soula guardrail: separate sensory following from capture of an endogenous oscillator",
           8.6, GREY, w=11, h=0.3)
    cards = [
        ("1. Steady-state response\n≠ entrainment",
         "A 40 Hz line can be\nsensory following.\nFlicker resonance +\nattention can amplify it.\nRefs: Herrmann 2001;\nTiitinen et al. 1993",
         "not native-gamma capture"),
        ("2. Soula/Buzsáki 2023",
         "40 Hz visual flicker:\ntrain-onset response\nV1 20.1% (92/458)\nEC 1.9%; CA1 0.2%\ncycle-locking\nV1 28.8%; CA1 6.7%; EC 3.3%",
         "visual cortex ≫ hippocampus / EC"),
        ("3. Independent caution",
         "Human MEG: rapid flicker\nproduced a strong visual\nresponse, but endogenous\ngamma and flicker response\ncoexisted rather than locking.",
         "independent caution"),
        ("4. What haptics must prove",
         "Record vibration phase;\nshow spikes + LFP lock;\nshow native rhythm is\npulled to the drive;\nrule out pickup; test\nspecificity + carryover.",
         "our next-study bar"),
    ]
    w, h, gap, y = 2.55, 2.35, 0.22, 0.82
    xs = [0.27 + i * (w + gap) for i in range(4)]
    for i, ((head, body, tag), x) in enumerate(zip(cards, xs), start=1):
        edge = RED if i in (2, 4) else GREY
        p.box(x, y, w, h, fc=FOG_LIGHT, ec=edge, lw=2.0)
        p.text(x + w / 2, y + h - 0.34, head, 8.9, RED if i == 2 else NAVY,
               bold=True, w=w - 0.2, h=0.52)
        p.text(x + w / 2, y + 1.18, body, 6.85, STONE_DARK, bold=True, w=w - 0.22, h=1.35)
        p.text(x + w / 2, y + 0.24, tag, 6.85, RED, bold=True, w=w - 0.2, h=0.28)
    for i in range(3):
        p.arrow(xs[i] + w + 0.02, y + h / 2, xs[i + 1] - 0.02, y + h / 2, color=GREY, lw=1.8)
    p.text(5.75, 0.28,
           "language discipline: response / target engagement first; entrainment only with native-rhythm phase evidence",
           9.1, RED, bold=True, w=11, h=0.35)


def driving_clearance(p: Pen):
    p.text(5.75, 3.65, "Preclinical proof: clearance can be driven non-invasively", 14, NAVY, bold=True, w=11)
    p.text(5.75, 3.35, "External stimulation examples in animals — still not human clinical proof", 9.3, GREY, w=11)
    p.box(0.4, 2.05, 3.0, 0.95, fc=FOG_LIGHT, ec=NAVY, lw=2.1)
    p.text(1.9, 2.525, "40 Hz light + sound\n(GENUS)", 11, NAVY, bold=True, w=2.9, h=0.9)
    p.arrow(3.45, 2.525, 4.3, 2.525, color=GREY, lw=2.4)
    p.box(4.35, 2.05, 6.75, 0.95, fc=FOG_LIGHT, ec=STONE_DARK, lw=2.1)
    p.text(7.725, 2.525, "mice: ~4× CSF influx · ~30% less amyloid · AQP4-dependent\n(Murdock 2024, Nature)", 10.5, STONE_DARK, bold=True, w=6.5, h=0.9)
    p.box(0.4, 0.75, 3.0, 0.95, fc=FOG_LIGHT, ec=STONE, lw=2.1)
    p.text(1.9, 1.225, "Focused ultrasound\n(Stanford)", 11, STONE_DARK, bold=True, w=2.9, h=0.9)
    p.arrow(3.45, 1.225, 4.3, 1.225, color=GREY, lw=2.4)
    p.box(4.35, 0.75, 6.75, 0.95, fc=FOG_LIGHT, ec=STONE, lw=2.1)
    p.text(7.725, 1.225, "rodents: ↑ CSF influx · faster debris clearance\n(Aryal 2022 · Azadian 2025, Nat Biotech)", 10.5, STONE_DARK, bold=True, w=6.5, h=0.9)
    p.text(5.75, 0.27, "open question: can VIBROTACTILE stimulation enhance glymphatic clearance?", 10.7, RED, bold=True, w=11)


def landscape(p: Pen):
    p.text(5.875, 4.02, "Approaches to Alzheimer's and the open frontier", 14, NAVY, bold=True, w=11.5)
    items = [
        ("Pharmacology",
         [("method: anti-amyloid antibodies", True), ("lecanemab: van Dyck NEJM 2023", False),
          ("donanemab: Sims JAMA 2023", False), ("impact: ~27-36% slower decline", True),
          ("but BBB + ARIA/MRI burden", False)], GREY),
        ("Lifestyle / multidomain",
         [("method: exercise/diet/cognitive/social", True), ("FINGER: Ngandu Lancet 2015", False),
          ("U.S. POINTER: Baker JAMA 2025", False), ("Stanford CT: Gozdas/Hosseini", False),
          ("Transl Psychiatry 2024", False), ("impact: +0.029 SD/yr; CT d~0.7", True),
          ("scope differs; MIND NEJM 2023 null", False)], GREEN),
        ("Non-invasive neuromodulation",
         [("method: rTMS / gamma tACS / 40 Hz", True), ("rTMS: Koch Alz Res Ther 2025", False),
          ("tACS: Cantoni JAMA Netw Open 2025", False), ("GENUS: Chan PLOS ONE 2022", False),
          ("impact: early / small trials", True)], RED),
    ]
    w, gap, y, h = 3.4, 0.5, 0.62, 2.85
    xs = [0.4 + k * (w + gap) for k in range(3)]
    for (title, lines, col), x in zip(items, xs):
        p.box(x, y, w, h, fc=LIGHT, ec=col, lw=2.6 if col == RED else 1.5)
        p.text(x + w / 2, y + h - 0.42, title, 12.3, col, bold=True, w=w)
        p.textblock(x + w / 2, y + h * 0.56, lines, 7.9, NAVY if col != GREY else GREY, w=w - 0.2, h=1.6)
    p.arrow(xs[2] + w / 2, 0.23, xs[2] + w / 2, y + 0.03, color=RED, lw=2.2)


def csf_wave_evidence(p: Pen):
    p.text(5.75, 4.28, "Neural waves can organize CSF perfusion and clearance", 14, NAVY, bold=True, w=11)
    p.text(5.75, 4.02, "Causal anchor plus converging sleep, vascular, and neuromodulator evidence", 9.6, GREY, w=11)
    # left causal panel
    p.box(0.35, 0.78, 5.15, 2.95, fc=FOG_LIGHT, ec=RED, lw=2.4)
    p.text(2.93, 3.48, "causal anchor: Jiang-Xie et al., Nature 2024", 10.6, RED, bold=True, w=5.0)
    T = [i * 2 * math.pi / 40 for i in range(41)]
    for k in range(3):
        p.curve([(0.85 + 1.05 * t / (2 * math.pi) + k * 0.18, 2.55 + 0.18 * math.sin(3 * t) - k * 0.32) for t in T],
                color=NAVY, lw=1.8)
    p.text(1.45, 2.92, "synchronized neurons", 9.4, NAVY, bold=True, w=2.6, h=0.35, anchor="bottom")
    p.arrow(2.1, 2.35, 2.82, 2.35, color=GREY, lw=2.1)
    p.curve([(2.95, 2.35), (3.25, 2.70), (3.55, 2.00), (3.85, 2.35)], color=POPPY, lw=2.4)
    p.text(3.4, 1.82, "ionic waves\nin tissue", 9.2, NAVY, w=1.6, h=0.5)
    p.arrow(3.98, 2.35, 4.75, 2.35, color=POPPY, lw=2.4)
    p.text(4.73, 2.72, "CSF → ISF\nperfusion", 9.2, POPPY, bold=True, w=1.7, h=0.5)
    p.text(2.9, 1.30, "method: flatten waves (chemogenetic) ↓ clearance\nsynthesize waves (transcranial optogenetic) ↑ perfusion",
           8.5, GREY, w=5.0, h=0.55)
    # right converging panel
    p.box(5.8, 0.78, 5.35, 2.95, fc=FOG_LIGHT, ec=NAVY, lw=1.7)
    p.text(8.48, 3.48, "not alone: converging support", 10.6, NAVY, bold=True, w=5.2)
    rows = [("human sleep", "Fultz 2019, Science", "slow waves + BOLD + CSF oscillate together", POPPY),
            ("sleep/glymphatic state", "Hablitz 2019; Hauglund 2025", "delta / norepinephrine / vasomotion link to influx", POPPY),
            ("neurovascular control", "Chuang 2025; Broggini 2024", "cholinergic and vascular waves regulate flow/perfusion", POPPY),
            ("external rhythm example", "Murdock 2024, Nature", "40 Hz light+sound drives glymphatic Aβ clearance in mice", RED)]
    for i, (label, cite, point, col) in enumerate(rows):
        yy = 3.12 - i * 0.62
        p.text(6.1, yy + 0.15, label, 8.5, NAVY, bold=True, align="left", w=4.7, h=0.28)
        p.text(6.1, yy - 0.02, cite, 8.0, col, align="left", w=4.7, h=0.28)
        p.text(6.1, yy - 0.20, point, 7.5, GREY, align="left", w=4.7, h=0.28)
    p.text(5.75, 0.30, "take-home: brain rhythms and state can gate clearance biology — motivating a vibrotactile rhythm test",
           9.8, RED, bold=True, w=11)


def glymphatic_measurement_methods(p: Pen):
    p.text(5.75, 4.22,
           "Measurement strength: human CSF motion -> direct mouse tracer flow -> causal 40 Hz clearance test",
           9.7, GREY, bold=True, w=11)

    cards = [
        ("Human CSF motion",
         "Fultz 2019 Science;\nWilliams 2023 PLOS Biol",
         "EFFECT:\nFultz +5.5 dB CSF power;\nWilliams +6-10% evoked CSF signal",
         "MEASURED:\nfast EEG/fMRI; BOLD and\nventricular CSF inflow waves",
         "TAKEAWAY:\nbrain state / sensory-evoked\nhemodynamics can move CSF",
         "LIMIT:\nnoninvasive but indirect;\nCSF motion ≠ solute clearance",
         LAGUNITA),
        ("Neural waves drive flow",
         "Jiang-Xie / Kipnis\nNature 2024",
         "EFFECT:\ntracer flow sleep/wake ~1.5-2x;\nsilencing impairs entry/clearance",
         "MEASURED:\nionic waves + CSF tracers;\nflatten waves ↓ clearance;\nsynthetic waves ↑ perfusion",
         "TAKEAWAY:\nneuronal dynamics can direct\nCSF-to-ISF perfusion",
         "LIMIT:\nmouse/preclinical;\nnot AD-specific or 40 Hz-specific",
         PALO_ALTO),
        ("40 Hz amyloid clearance",
         "Murdock / Tsai\nNature 2024",
         "EFFECT:\nCSF influx ~4x; ISF efflux ~1.5x;\ncortical Aβ ↓~30%",
         "MEASURED:\ncisterna magna tracer influx,\nISF efflux, Aβ in cervical nodes,\nAQP4 blockade",
         "TAKEAWAY:\n40 Hz light+sound increased\nglymphatic flow and Aβ clearance",
         "LIMIT:\nclosest to our story, but still\nmouse; audiovisual not haptic",
         RED),
    ]
    x0s = [0.35, 4.10, 7.85]
    w, h, y = 3.30, 3.22, 0.78
    for i, (head, cite, effect, measured, takeaway, limit, col) in enumerate(cards):
        x = x0s[i]
        p.box(x, y, w, h, fc=FOG_LIGHT, ec=col, lw=2.1 if i == 2 else 1.7, rounded=False)
        p.text(x + w / 2, y + h - 0.26, head, 10.0, col, bold=True, w=w - 0.2, h=0.33)
        p.text(x + w / 2, y + h - 0.62, cite, 7.5, NAVY, bold=True, w=w - 0.25, h=0.42)
        p.text(x + 0.24, y + h - 1.04, effect, 5.95, col, bold=True, align="left", w=w - 0.46, h=0.50)
        p.line(x + 0.22, y + h - 1.34, x + w - 0.22, y + h - 1.34, color=STONE_DARK, lw=0.65)
        p.text(x + 0.24, y + h - 1.70, measured, 6.55, STONE_DARK, align="left", w=w - 0.46, h=0.62)
        p.text(x + 0.24, y + h - 2.38, takeaway, 6.55, NAVY, bold=True, align="left", w=w - 0.46, h=0.58)
        p.text(x + 0.24, y + 0.25, limit, 6.05, GREY, align="left", w=w - 0.46, h=0.50)

    for i in range(2):
        p.arrow(x0s[i] + w + 0.06, y + h / 2, x0s[i + 1] - 0.06, y + h / 2,
                color=STONE, lw=1.7)

    p.box(0.80, 0.04, 9.90, 0.42, fc="#FFF7F4", ec=RED, lw=1.2, rounded=False)
    p.text(5.75, 0.27,
           "Program payoff: evaluate haptics with flow / clearance biomarkers, not entrainment language alone.",
           8.5, RED, bold=True, w=9.6, h=0.36)


def vibrotactile_anchor_suk(p: Pen):
    p.text(5.75, 4.02, "Suk 2023 is the haptic precedent — and it defines the missing measurement",
           13.4, NAVY, bold=True, w=10.8)
    p.text(5.75, 3.74,
           "whole-body 40 Hz vibration in AD-model mice; activity/pathology readouts, not direct electrophysiology",
           8.6, GREY, w=10.8)

    cards = [
        ("1. Stimulus route",
         "40 Hz whole-body\nvibration using a\nlarge speaker / platform",
         "not a wearable\nfingertip device",
         GREY),
        ("2. Brain readout",
         "immunostaining:\nincreased c-Fos activity\nmarkers in superficial\nSSp + MOp cortex",
         "activity marker,\nnot phase timing",
         PALO_ALTO),
        ("3. What was not tested",
         "no probes; no LFP phase-locking;\nno single-unit entrainment;\nno hippocampus or entorhinal\nrecordings",
         "the gap our study targets",
         RED),
    ]
    x0s = [0.35, 4.10, 7.85]
    w, h, y = 3.30, 2.55, 0.80
    for i, (head, body, tag, col) in enumerate(cards):
        x = x0s[i]
        p.box(x, y, w, h, fc=FOG_LIGHT, ec=col, lw=2.2 if i == 2 else 1.7, rounded=False)
        p.text(x + w / 2, y + h - 0.30, head, 9.7, col if i else NAVY, bold=True, w=w - 0.2, h=0.34)
        p.text(x + w / 2, y + 1.33, body, 8.1, NAVY, bold=True, w=w - 0.34, h=1.30)
        p.text(x + w / 2, y + 0.30, tag, 7.4, RED if i == 2 else GREY, bold=True, w=w - 0.35, h=0.46)

    for i in range(2):
        p.arrow(x0s[i] + w + 0.06, y + h / 2, x0s[i + 1] - 0.06, y + h / 2,
                color=STONE, lw=1.8)

    p.box(0.86, 0.18, 9.78, 0.42, fc="#FFF7F4", ec=RED, lw=1.2, rounded=False)
    p.text(5.75, 0.39,
           "Bridge to us: keep the haptic promise, but prove direct medial-temporal target engagement with spikes + LFP.",
           8.5, RED, bold=True, w=9.4, h=0.36)


def frequency_menu(p: Pen):
    p.text(5.75, 4.42, "Frequency menu: each condition answers a different question",
           13.5, NAVY, bold=True, w=10.8)
    p.text(5.75, 4.15, "Not one magic number: test literature anchor, our strongest response, and controls",
           8.8, GREY, w=10.8)

    rows = [
        ("~0.05-1 Hz\nslow / sleep-state",
         "Fultz 2019;\nJiang-Xie/Kipnis 2024",
         "CSF waves during sleep; synthetic 1 Hz neural waves\npotentiate CSF-to-ISF perfusion",
         "clearance biology may be slow-state,\nnot gamma-only",
         PALO_ALTO),
        ("5 / 10 Hz\nlow flutter controls",
         "Mountcastle/Romo;\nHayashi 2018",
         "classic tactile flutter range; mouse S1 responds\nfrequency-selectively to vibration",
         "asks whether any periodic touch\ndrives the brain",
         GREY),
        ("20 / 26 Hz\nmid controls",
         "Romo/Hayashi;\nour Dec 3/4",
         "mid-flutter / beta comparator; 26 Hz gave large\nLFP transients but not the clean unit peak",
         "protects against cherry-picking\n40/50 Hz",
         STONE_DARK),
        ("40 Hz\nAD anchor",
         "Tsai/GENUS;\nMurdock 2024; Suk 2023",
         "AD-model pathology/glymphatic anchor; whole-body\nvibration precedent, but no direct entrainment test",
         "required literature comparator",
         RED),
        ("50 Hz\nour candidate",
         "Romo/Hayashi;\nour Dec 4 spikes",
         "upper flutter; strongest clean single-unit firing-rate\nmodulation in dHPC + LEC",
         "test with phase reference +\nartifact controls",
         LAGUNITA),
        ("250 Hz pattern\noptional CR",
         "Tass 2017;\nPfeifer/Tass 2021",
         "fingertip vibrotactile coordinated-reset uses high-frequency\nbursts with a slow multi-site pattern",
         "patterning idea;\nnot AD-gamma evidence",
         POPPY),
    ]
    x0, y0 = 0.38, 3.55
    widths = [1.55, 2.00, 4.05, 2.75]
    row_h = 0.49
    xs = [x0]
    for ww in widths[:-1]:
        xs.append(xs[-1] + ww)
    total_w = sum(widths)
    p.box(x0, y0, total_w, 0.35, fc=FOG_LIGHT, ec=STONE, lw=0.7, rounded=False)
    for label, x, ww in zip(["Frequency", "Reference paper(s)", "Finding / prompt", "Why include it"], xs, widths):
        p.text(x + 0.07, y0 + 0.18, label, 7.2, NAVY, bold=True, align="left", w=ww - 0.12, h=0.22)
    y = y0 - row_h
    for idx, (freq, refs, finding, why, col) in enumerate(rows):
        p.box(x0, y, total_w, row_h, fc="#FFFFFF" if idx % 2 else "#FAFAFA", ec=STONE, lw=0.4, rounded=False)
        p.box(x0, y, 0.08, row_h, fc=col, ec=col, lw=0, rounded=False)
        p.text(xs[0] + 0.15, y + row_h / 2, freq, 6.6, col, bold=True, align="left", w=widths[0] - 0.2, h=row_h - 0.04)
        p.text(xs[1] + 0.08, y + row_h / 2, refs, 6.1, NAVY, bold=True, align="left", w=widths[1] - 0.16, h=row_h - 0.04)
        p.text(xs[2] + 0.08, y + row_h / 2, finding, 5.85, STONE_DARK, align="left", w=widths[2] - 0.16, h=row_h - 0.04)
        p.text(xs[3] + 0.08, y + row_h / 2, why, 5.85, NAVY, bold=True, align="left", w=widths[3] - 0.16, h=row_h - 0.04)
        y -= row_h
    for x in xs[1:]:
        p.line(x, y0 + 0.35, x, y0 - row_h * len(rows), color=STONE, lw=0.55)

    p.box(0.70, 0.12, 10.10, 0.38, fc="#FFF7F4", ec=RED, lw=1.2, rounded=False)
    p.text(5.75, 0.31,
           "Core panel: 5, 10, 20, 26, 40, 50 Hz; slow/patterned = optional clearance or CR extensions.",
           7.55, RED, bold=True, w=9.85, h=0.34)


def glymphatic(p: Pen):
    p.oval(5.75, 2.55, 4.85, 1.525, ec="#D4D1D1", lw=1.9)
    p.line(5.75, 1.10, 5.75, 4.00, color="#D4D1D1", lw=1.6, dash="dash")
    p.oval(5.75, 2.48, 2.20, 0.61, ec=CARDINAL, lw=2.4)
    p.oval(3.42, 2.56, 0.34, 0.23, ec=CARDINAL, lw=1.8)
    p.oval(7.88, 2.39, 0.275, 0.31, ec=CARDINAL, lw=1.8)
    # left: more flow
    p.text(3.25, 3.72, "NREM sleep", 15.5, NAVY, bold=True, w=3.0)
    p.text(3.25, 3.40, "more CSF flow", 10.3, GREY, w=3.0, h=0.3)
    p.arrow(1.95, 3.28, 3.12, 2.82, color=STONE_DARK, lw=4.8)
    p.arrow(4.80, 3.22, 3.86, 2.82, color=STONE_DARK, lw=4.2)
    p.arrow(3.80, 2.18, 2.95, 1.22, color=STONE_DARK, lw=5.0)
    p.text(2.15, 1.32, "more\nwaste out", 10.0, CARDINAL, bold=True, w=1.6, h=0.5)
    p.text(3.16, 1.00, "toward cervical\nlymph nodes", 7.6, GREY, w=2.0, h=0.4)
    # right: less flow
    p.text(8.25, 3.72, "Awake", 15.5, NAVY, bold=True, w=3.0)
    p.text(8.25, 3.40, "less CSF flow", 10.3, GREY, w=3.0, h=0.3)
    p.arrow(9.70, 3.25, 8.58, 2.82, color=STONE, lw=2.2)
    p.arrow(8.28, 2.15, 8.78, 1.35, color=STONE, lw=2.1)
    p.text(9.30, 1.32, "clearance\ncontinues", 8.9, GREY, w=1.8, h=0.4)
    rng = [(2.9, 2.3), (3.6, 2.6), (4.2, 2.4), (3.2, 2.7), (4.4, 2.2), (3.9, 2.5)]
    for (dx, dy) in rng:
        p.oval(dx, dy, 0.05, 0.05, fc=CARDINAL)
    for dx in [6.9, 7.3, 7.7, 8.1, 8.5, 8.9, 7.1, 7.5, 7.9, 8.3, 8.7, 7.6]:
        p.oval(dx, 2.2 + (dx % 0.4), 0.05, 0.05, fc=CARDINAL)
    p.text(9.18, 2.70, "amyloid-β / tau", 9.4, CARDINAL, align="left", w=2.2, h=0.3)
    p.text(5.75, 0.48, "Direct tracer evidence is strongest in mice/rodents; human studies show coupled sleep, blood-flow, and CSF rhythms.",
           7.9, GREY, w=11, h=0.3)
    p.text(5.75, 0.24, "Refs: Xie 2013; Hablitz 2019; Fultz 2019; Hauglund 2025.", 7.3, GREY, w=11, h=0.3)


# stem -> (xmax, ymax, renderer)
NATIVE = {
    "whitespace": (11, 4.2, whitespace),
    "vision_chain": (12, 3.4, vision_chain),
    "genus_evidence_chain": (11.5, 4.2, genus_evidence_chain),
    "response_entrainment_guardrail": (11.5, 4.2, response_entrainment_guardrail),
    "driving_clearance": (11.5, 3.9, driving_clearance),
    "landscape": (11.75, 4.25, landscape),
    "csf_wave_evidence": (11.5, 4.55, csf_wave_evidence),
    "glymphatic_measurement_methods": (11.5, 4.55, glymphatic_measurement_methods),
    "vibrotactile_anchor_suk": (11.5, 4.25, vibrotactile_anchor_suk),
    "frequency_menu": (11.5, 4.65, frequency_menu),
    "glymphatic": (11.5, 4.1, glymphatic),
}


def render(s, stem, bx, by, bw, bh):
    """Draw the native diagram for `stem` into the slide box (inches). Returns True if drawn."""
    if stem not in NATIVE:
        return False
    xmax, ymax, fn = NATIVE[stem]
    fn(Pen(s, xmax, ymax, bx, by, bw, bh))
    return True
